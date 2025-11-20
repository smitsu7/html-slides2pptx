#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Browser-free HTML -> PPTX converter with virtual rendering.

Key ideas
---------
- StyleParser: lightweight CSS cascade (inline + <style>) without cssutils.
- LayoutEngine: approximates block/float/flex/grid flow to preserve card and
  multi-column layouts and paginates when vertical space overflows.
- StyleMapping: maps CSS-like colors/borders/typography to python-pptx shapes.
- ChartBuilder: parses Chart.js-style scripts to rebuild editable PPT charts
  instead of screenshot images.

Usage:
    python html2pptx_light_3.py input.html output.pptx

Requirements:
    pip install beautifulsoup4 python-pptx pillow lxml
"""

from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

try:
    from PIL import Image, ImageColor, ImageDraw, ImageFont
except ImportError:  # Pillow is optional but recommended
    Image = None
    ImageColor = None
    ImageDraw = None
    ImageFont = None

# --------------------------------------------------------------------------- #
# Constants
# --------------------------------------------------------------------------- #

SLIDE_REF_WIDTH = 1920.0
SLIDE_REF_HEIGHT = 1080.0
PX_PER_INCH = 96.0
# Keep a tiny safety margin but primarily respect HTML padding.
BASE_MARGIN_X = 12.0
BASE_MARGIN_Y = 12.0
FLOW_GAP_Y = 28.0
FLOW_GAP_X = 24.0
DEFAULT_FONT_SIZE = 28.0
DEFAULT_FONT_FAMILY = "Arial"
TAG_FONT_SIZE = {"h1": 64.0, "h2": 48.0, "h3": 36.0, "h4": 30.0, "h5": 28.0, "h6": 26.0, "li": 26.0, "p": 28.0}
INHERITED_PROPS = {"color", "font-size", "font-family", "font-weight", "font-style", "text-align", "line-height"}

COLOR_KEYWORDS = {
    "white": (255, 255, 255),
    "black": (0, 0, 0),
    "red": (255, 0, 0),
    "green": (0, 128, 0),
    "lime": (0, 255, 0),
    "blue": (0, 0, 255),
    "yellow": (255, 255, 0),
    "orange": (255, 165, 0),
    "purple": (128, 0, 128),
    "gray": (128, 128, 128),
    "grey": (128, 128, 128),
    "silver": (192, 192, 192),
    "navy": (0, 0, 128),
    "teal": (0, 128, 128),
    "cyan": (0, 255, 255),
    "aqua": (0, 255, 255),
    "magenta": (255, 0, 255),
    "fuchsia": (255, 0, 255),
    "maroon": (128, 0, 0),
    "olive": (128, 128, 0),
}

# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #


def px(value: float) -> float:
    return float(value)


def parse_length(value: Optional[str], reference: Optional[float] = None) -> Optional[float]:
    if value is None:
        return None
    value = str(value).strip()
    if not value:
        return None
    lowered = value.lower()
    if lowered.endswith("px"):
        try:
            return float(lowered[:-2])
        except ValueError:
            return None
    if lowered.endswith("pt"):  # 1pt = 1.333px approximate
        try:
            return float(lowered[:-2]) * 1.333
        except ValueError:
            return None
    if lowered.endswith("%") and reference is not None:
        try:
            return float(lowered[:-1]) / 100.0 * reference
        except ValueError:
            return None
    try:
        return float(lowered)
    except ValueError:
        return None


def parse_box(value: Optional[str], fallback: float = 0.0) -> Tuple[float, float, float, float]:
    if not value:
        return (fallback, fallback, fallback, fallback)
    tokens = [t for t in value.replace(",", " ").split() if t]
    nums = [parse_length(t) if parse_length(t) is not None else fallback for t in tokens]
    if len(nums) == 1:
        return (nums[0], nums[0], nums[0], nums[0])
    if len(nums) == 2:
        return (nums[0], nums[1], nums[0], nums[1])
    if len(nums) == 3:
        return (nums[0], nums[1], nums[2], nums[1])
    t, r, b, l = (nums + [fallback] * 4)[:4]
    return (t, r, b, l)


def parse_color(value: Optional[str]) -> Optional[RGBColor]:
    if not value:
        return None
    raw = value.strip()
    if not raw or raw.lower() == "transparent":
        return None
    if raw.startswith("#"):
        hex_val = raw[1:]
        if len(hex_val) == 3:
            hex_val = "".join([c * 2 for c in hex_val])
        if len(hex_val) == 6:
            try:
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                return RGBColor(r, g, b)
            except ValueError:
                return None
    if raw.lower() in COLOR_KEYWORDS:
        r, g, b = COLOR_KEYWORDS[raw.lower()]
        return RGBColor(r, g, b)
    if raw.startswith("rgb"):
        try:
            inside = raw[raw.find("(") + 1 : raw.find(")")]
            parts = [p.strip() for p in inside.split(",")]
            r, g, b = [int(float(p)) for p in parts[:3]]
            return RGBColor(r, g, b)
        except Exception:
            return None
    if ImageColor:
        try:
            r, g, b = ImageColor.getrgb(raw)
            return RGBColor(r, g, b)
        except Exception:
            return None
    return None


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def pt_from_px(px_val: float) -> Pt:
    return Pt(px_val * 72.0 / PX_PER_INCH)


# --------------------------------------------------------------------------- #
# CSS parsing without cssutils
# --------------------------------------------------------------------------- #


@dataclass
class CSSRule:
    selector: str
    selector_chain: List[str]
    props: Dict[str, str]
    specificity: Tuple[int, int, int]
    order: int


class StyleParser:
    TOKEN_RE = re.compile(r"([#.]?[\w-]+|\*)")
    RULE_RE = re.compile(r"([^{}]+)\{([^{}]+)\}")
    COMMENT_RE = re.compile(r"/\*.*?\*/", re.S)

    def __init__(self, soup: BeautifulSoup):
        self.soup = soup
        self.rules: List[CSSRule] = []
        self._collect_rules()

    def _collect_rules(self):
        order = 0
        for style_tag in self.soup.find_all("style"):
            css_text = style_tag.string or ""
            cleaned = self.COMMENT_RE.sub("", css_text)
            for sel_text, body in self.RULE_RE.findall(cleaned):
                selectors = [s.strip() for s in sel_text.split(",") if s.strip()]
                props = self._parse_declarations(body)
                for selector in selectors:
                    chain = [part for part in selector.split() if part]
                    spec = self._specificity(chain)
                    self.rules.append(CSSRule(selector, chain, props, spec, order))
                    order += 1

    def _parse_declarations(self, body: str) -> Dict[str, str]:
        props: Dict[str, str] = {}
        for decl in body.split(";"):
            if ":" not in decl:
                continue
            key, val = decl.split(":", 1)
            key = key.strip().lower()
            val = val.strip()
            if not key:
                continue
            props[key] = val
        return props

    def _specificity(self, chain: List[str]) -> Tuple[int, int, int]:
        ids = sum(1 for token in chain for part in self.TOKEN_RE.findall(token) if part.startswith("#"))
        classes = sum(1 for token in chain for part in self.TOKEN_RE.findall(token) if part.startswith("."))
        tags = sum(1 for token in chain for part in self.TOKEN_RE.findall(token) if not part.startswith(("#", ".")) and part != "*")
        return (ids, classes, tags)

    def _match_component(self, element: Tag, token: str) -> bool:
        tag = None
        id_val = None
        classes: List[str] = []
        for part in self.TOKEN_RE.findall(token):
            if part == "*":
                continue
            if part.startswith("#"):
                id_val = part[1:]
            elif part.startswith("."):
                classes.append(part[1:])
            else:
                tag = part.lower()
        if tag and element.name != tag:
            return False
        if id_val and element.get("id") != id_val:
            return False
        elem_classes = set(element.get("class", []))
        if classes and not all(cls in elem_classes for cls in classes):
            return False
        return True

    def _matches(self, element: Tag, chain: List[str]) -> bool:
        if not chain:
            return False
        idx = len(chain) - 1
        current: Optional[Tag] = element
        while current and idx >= 0:
            if self._match_component(current, chain[idx]):
                idx -= 1
            current = current.parent if isinstance(current.parent, Tag) else None
        return idx < 0

    def compute(self, element: Tag, parent_style: Optional[Dict[str, str]] = None) -> Dict[str, str]:
        style: Dict[str, str] = {}
        if parent_style:
            for key in INHERITED_PROPS:
                if key in parent_style:
                    style[key] = parent_style[key]
        tag_defaults = {}
        if element.name in TAG_FONT_SIZE:
            tag_defaults["font-size"] = f"{TAG_FONT_SIZE[element.name]}px"
            if element.name.startswith("h"):
                tag_defaults["font-weight"] = "700"
        style.update(tag_defaults)

        applicable = [rule for rule in self.rules if self._matches(element, rule.selector_chain)]
        applicable.sort(key=lambda r: (r.specificity, r.order))
        for rule in applicable:
            style.update(rule.props)

        inline = element.get("style")
        if inline:
            style.update(self._parse_declarations(inline))

        # Attribute-based hints
        if element.name in {"ul", "ol"} and "list-style-type" not in style:
            style["list-style-type"] = "disc" if element.name == "ul" else "decimal"
        if "display" not in style:
            style["display"] = "block" if element.name not in {"span", "b", "strong", "em", "i"} else "inline"
        if "color" not in style and element.name in {"a"}:
            style["color"] = "#1155cc"
        return style


# --------------------------------------------------------------------------- #
# Data classes used for rendering
# --------------------------------------------------------------------------- #


@dataclass
class TextRun:
    text: str
    font_size: Optional[float] = None
    bold: bool = False
    italic: bool = False
    color: Optional[str] = None


@dataclass
class LayoutFrame:
    left: float = 0.0
    top: float = 0.0
    width: float = 0.0
    height: float = 0.0


@dataclass
class ChartSeries:
    name: str
    values: List[float]
    color: Optional[str] = None


@dataclass
class ChartSpec:
    chart_type: str
    labels: List[str]
    series: List[ChartSeries]
    target_ids: List[str] = field(default_factory=list)


@dataclass
class Block:
    kind: str
    style: Dict[str, str]
    frame: LayoutFrame = field(default_factory=LayoutFrame)
    text: str = ""
    runs: List[TextRun] = field(default_factory=list)
    items: List[str] = field(default_factory=list)
    table: List[List[str]] = field(default_factory=list)
    cell_styles: List[List[Dict[str, str]]] = field(default_factory=list)
    image_path: Optional[Path] = None
    chart: Optional[ChartSpec] = None
    children: List["Block"] = field(default_factory=list)


# --------------------------------------------------------------------------- #
# Text measurement helpers
# --------------------------------------------------------------------------- #


class TextMeasurer:
    def __init__(self):
        if ImageDraw:
            self.canvas = Image.new("RGB", (1, 1))
            self.drawer = ImageDraw.Draw(self.canvas)
        else:
            self.canvas = None
            self.drawer = None
        self.font_cache: Dict[Tuple[str, int, str], Any] = {}

    def _font(self, family: str, size: float, weight: str):
        key = (family, int(size), weight or "normal")
        if key in self.font_cache:
            return self.font_cache[key]
        font = None
        if ImageFont:
            candidates = ["Arial.ttf", "/System/Library/Fonts/Supplemental/Arial.ttf"]
            for cand in candidates:
                try:
                    font = ImageFont.truetype(cand, int(size))
                    break
                except OSError:
                    continue
            if font is None:
                try:
                    font = ImageFont.load_default()
                except Exception:
                    font = None
        self.font_cache[key] = font
        return font

    def measure_block(self, text: str, style: Dict[str, str], width: float) -> Tuple[float, float]:
        if not text:
            return (width, DEFAULT_FONT_SIZE * 1.2)
        font_size = parse_length(style.get("font-size")) or DEFAULT_FONT_SIZE
        family = style.get("font-family", DEFAULT_FONT_FAMILY)
        weight = style.get("font-weight", "normal")
        if not self.drawer or not self.canvas:
            # Rough estimation fallback
            lines = max(1, len(text.splitlines()))
            return (width, lines * font_size * 1.25 + 10)

        font = self._font(family, font_size, weight)
        words = text.replace("\xa0", " ").split()
        space_w = self.drawer.textbbox((0, 0), " ", font=font)[2]
        lines: List[str] = []
        current = ""
        current_w = 0
        max_w = 0
        for word in words:
            w = self.drawer.textbbox((0, 0), word, font=font)[2]
            if current and current_w + space_w + w > width:
                lines.append(current)
                max_w = max(max_w, current_w)
                current = word
                current_w = w
            else:
                if current:
                    current += " "
                    current_w += space_w + w
                else:
                    current = word
                    current_w = w
        if current:
            lines.append(current)
            max_w = max(max_w, current_w)
        # line height
        lh_box = self.drawer.textbbox((0, 0), "Ag", font=font)
        line_height = (lh_box[3] - lh_box[1]) * 1.2
        height = max(1, len(lines)) * line_height + 10
        return (max_w, height)


# --------------------------------------------------------------------------- #
# Inline text extraction
# --------------------------------------------------------------------------- #


class TextRunExtractor:
    def __init__(self, style_parser: StyleParser):
        self.styles = style_parser

    def extract(self, element: Tag, base_style: Dict[str, str]) -> List[TextRun]:
        runs: List[TextRun] = []

        def walk(node, inherited: Dict[str, str]):
            if isinstance(node, NavigableString):
                raw = str(node).replace("\xa0", " ")
                if raw.strip():
                    runs.append(self._run(raw, inherited))
                return
            if not isinstance(node, Tag):
                return
            if node.name in {"script", "style"}:
                return
            local_style = self.styles.compute(node, inherited.copy())
            if node.name in {"strong", "b"}:
                local_style["font-weight"] = "bold"
            if node.name in {"em", "i"}:
                local_style["font-style"] = "italic"
            if node.name == "br":
                runs.append(self._run("\n", local_style))
                return
            for child in node.children:
                walk(child, local_style)

        walk(element, base_style.copy())
        merged: List[TextRun] = []
        for run in runs:
            if merged and self._can_merge(merged[-1], run):
                merged[-1].text += run.text
            else:
                merged.append(run)
        return merged

    def _run(self, text: str, style: Dict[str, str]) -> TextRun:
        return TextRun(
            text=text,
            font_size=parse_length(style.get("font-size")),
            bold=self._is_bold(style.get("font-weight")),
            italic=(style.get("font-style") == "italic"),
            color=style.get("color"),
        )

    def _is_bold(self, value: Optional[str]) -> bool:
        if not value:
            return False
        if value.isdigit():
            return int(value) >= 600
        return value.lower() in {"bold", "bolder", "700", "800", "900"}

    def _can_merge(self, a: TextRun, b: TextRun) -> bool:
        return a.font_size == b.font_size and a.bold == b.bold and a.italic == b.italic and a.color == b.color


# --------------------------------------------------------------------------- #
# DOM -> Block tree
# --------------------------------------------------------------------------- #


class BlockBuilder:
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp"}

    def __init__(self, style_parser: StyleParser, base_dir: Path, chart_specs: List[ChartSpec]):
        self.styles = style_parser
        self.base_dir = base_dir
        self.text_runs = TextRunExtractor(style_parser)
        self.chart_specs = chart_specs
        self.chart_lookup: Dict[str, ChartSpec] = {}
        for spec in chart_specs:
            for cid in spec.target_ids:
                self.chart_lookup[cid] = spec

    def build(self, container: Tag, parent_style: Optional[Dict[str, str]] = None) -> List[Block]:
        blocks: List[Block] = []
        parent_style = parent_style or self.styles.compute(container, None)
        for child in container.children:
            if isinstance(child, NavigableString):
                # Ignore stray text inside non-visible nodes (script/style/noscript)
                if isinstance(container, Tag) and container.name in {"script", "style", "noscript"}:
                    continue
                if normalize_text(str(child)):
                    pseudo = Tag(name="p")
                    pseudo.string = str(child)
                    pseudo["style"] = ""
                    block = self._text_block(pseudo, parent_style)
                    if block:
                        blocks.append(block)
                continue
            if isinstance(child, Tag):
                blk = self._dispatch(child, parent_style)
                if blk:
                    blocks.append(blk)
        return blocks

    def _dispatch(self, tag: Tag, parent_style: Dict[str, str]) -> Optional[Block]:
        if tag.name in {"script", "style"}:
            return None
        style = self.styles.compute(tag, parent_style)
        if tag.name in {"h1", "h2", "h3", "h4", "h5", "h6", "p", "span", "blockquote"}:
            return self._text_block(tag, style)
        if tag.name == "div":
            if not self._has_block_child(tag):
                return self._text_block(tag, style)
            # otherwise treat as container group
        if tag.name in {"ul", "ol"}:
            return self._list_block(tag, style)
        if tag.name == "table":
            return self._table_block(tag, style)
        if tag.name == "img":
            return self._image_block(tag, style)
        if tag.name == "canvas":
            return self._chart_block(tag, style)

        # Generic container
        children = self.build(tag, style)
        blk = Block("group", style=style, children=children)
        return blk if (children or self._has_shape(style)) else None

    def _text_block(self, tag: Tag, style: Dict[str, str]) -> Optional[Block]:
        text = normalize_text(tag.get_text(" ", strip=False))
        runs = self.text_runs.extract(tag, style)
        if not text and not runs:
            return None
        block = Block("text", style=style, text=text, runs=runs)
        return block

    def _list_block(self, tag: Tag, style: Dict[str, str]) -> Optional[Block]:
        items = [normalize_text(li.get_text(" ", strip=True)) for li in tag.find_all("li", recursive=False)]
        items = [i for i in items if i]
        if not items:
            return None
        return Block("list", style=style, items=items)

    def _table_block(self, tag: Tag, style: Dict[str, str]) -> Optional[Block]:
        rows: List[List[str]] = []
        cell_styles: List[List[Dict[str, str]]] = []
        for tr in tag.find_all("tr", recursive=False):
            row: List[str] = []
            row_styles: List[Dict[str, str]] = []
            for cell in tr.find_all(["td", "th"], recursive=False):
                row.append(normalize_text(cell.get_text(" ", strip=True)))
                row_styles.append(self.styles.compute(cell, style))
            if row:
                rows.append(row)
                cell_styles.append(row_styles)
        if not rows:
            return None
        return Block("table", style=style, table=rows, cell_styles=cell_styles)

    def _image_block(self, tag: Tag, style: Dict[str, str]) -> Optional[Block]:
        src = tag.get("src", "")
        img_path = self._resolve_image(src)
        return Block("image", style=style, image_path=img_path, text=tag.get("alt", ""))

    def _chart_block(self, tag: Tag, style: Dict[str, str]) -> Block:
        chart_spec = None
        cid = tag.get("id")
        if cid and cid in self.chart_lookup:
            chart_spec = self.chart_lookup[cid]
        # Carry explicit width/height attributes into style to guide layout.
        style = style.copy()
        if "width" not in style and tag.get("width"):
            style["width"] = f"{tag.get('width')}px"
        if "height" not in style and tag.get("height"):
            style["height"] = f"{tag.get('height')}px"
        return Block("chart", style=style, chart=chart_spec)

    def _resolve_image(self, src: str) -> Optional[Path]:
        if not src:
            return None
        p = Path(src)
        if p.is_absolute() and p.exists():
            return p
        candidate = (self.base_dir / src).resolve()
        if candidate.exists():
            return candidate
        return None

    def _has_shape(self, style: Dict[str, str]) -> bool:
        return bool(style.get("background") or style.get("background-color") or style.get("border") or style.get("border-color"))

    def _has_block_child(self, tag: Tag) -> bool:
        for child in tag.children:
            if isinstance(child, Tag) and child.name not in {"span", "b", "strong", "em", "i", "br"}:
                return True
        return False


# --------------------------------------------------------------------------- #
# Layout Engine with float/flex-aware flow and pagination
# --------------------------------------------------------------------------- #


class LayoutEngine:
    def __init__(self, width: float, height: float, measurer: TextMeasurer):
        self.width = width
        self.height = height
        self.measurer = measurer

    def layout_slide(self, blocks: List[Block], slide_style: Dict[str, str]) -> List[List[Block]]:
        pad_top, pad_right, pad_bottom, pad_left = parse_box(slide_style.get("padding"), 0)
        cursor_y = BASE_MARGIN_Y + pad_top
        usable_width = self.width - 2 * BASE_MARGIN_X - pad_left - pad_right
        flat: List[Block] = []
        for block in blocks:
            cursor_y = self._layout_block(block, BASE_MARGIN_X + pad_left, cursor_y, usable_width, flat)
            cursor_y += FLOW_GAP_Y
        return self._paginate(flat)

    def _layout_block(self, block: Block, left: float, top: float, width: float, collected: List[Block]) -> float:
        margin_top, margin_right, margin_bottom, margin_left = parse_box(block.style.get("margin"), 0)
        padding_top, padding_right, padding_bottom, padding_left = parse_box(block.style.get("padding"), 0)
        block_left = left + margin_left
        available_width = max(40.0, width - margin_left - margin_right)
        explicit_width = parse_length(block.style.get("width"), width)
        content_width = explicit_width if explicit_width else available_width

        if block.kind == "group":
            direction = block.style.get("display", "block")
            if "grid" in direction:
                self._layout_grid(block, block_left, top + margin_top, content_width, padding_top, padding_right, padding_bottom, padding_left, collected)
            elif "flex" in direction:
                self._layout_flex(block, block_left, top + margin_top, content_width, padding_top, padding_right, padding_bottom, padding_left, collected)
        else:
            self._layout_flow(block, block_left, top + margin_top, content_width, padding_top, padding_right, padding_bottom, padding_left, collected)
            # height already set inside helpers
            return block.frame.top + block.frame.height + margin_bottom

        # leaf blocks
        explicit_height = parse_length(block.style.get("height"), content_width)
        measured = self._measure_block(block, content_width - padding_left - padding_right)
        height = (explicit_height if explicit_height is not None else measured) + padding_top + padding_bottom
        block.frame = LayoutFrame(block_left, top + margin_top, content_width, height)
        collected.append(block)
        return block.frame.top + block.frame.height + margin_bottom

    def _layout_flow(
        self,
        block: Block,
        left: float,
        top: float,
        width: float,
        pad_top: float,
        pad_right: float,
        pad_bottom: float,
        pad_left: float,
        collected: List[Block],
    ):
        start_index = len(collected)
        inner_left = left + pad_left
        content_width = max(40.0, width - pad_left - pad_right)
        cursor_y = top + pad_top
        line_left = inner_left
        line_right = inner_left + content_width
        line_height = 0.0
        for child in block.children:
            float_dir = child.style.get("float", "").lower()
            if float_dir in {"left", "right"}:
                target_width = parse_length(child.style.get("width"), content_width) or max(content_width / 2.0, 120.0)
                if float_dir == "left":
                    if line_left + target_width > line_right and line_left > inner_left:
                        cursor_y += line_height + FLOW_GAP_Y
                        line_left = inner_left
                        line_right = inner_left + content_width
                        line_height = 0.0
                    self._layout_block(child, line_left, cursor_y, target_width, collected)
                    line_height = max(line_height, child.frame.height)
                    line_left += child.frame.width + FLOW_GAP_X
                else:  # right float
                    if line_right - target_width < line_left and line_right < inner_left + content_width:
                        cursor_y += line_height + FLOW_GAP_Y
                        line_left = inner_left
                        line_right = inner_left + content_width
                        line_height = 0.0
                    self._layout_block(child, line_right - target_width, cursor_y, target_width, collected)
                    line_height = max(line_height, child.frame.height)
                    line_right -= child.frame.width + FLOW_GAP_X
            else:
                if line_height:
                    cursor_y += line_height + FLOW_GAP_Y
                    line_left = inner_left
                    line_right = inner_left + content_width
                    line_height = 0.0
                cursor_y = self._layout_block(child, inner_left, cursor_y, content_width, collected)
                cursor_y += FLOW_GAP_Y
        total_height = max(
            pad_top + pad_bottom,
            (cursor_y - top) + line_height + pad_bottom if line_height else (cursor_y - top) + pad_bottom,
        )
        block.frame = LayoutFrame(left, top, width, total_height)
        collected.insert(start_index, block)

    def _layout_flex(
        self,
        block: Block,
        left: float,
        top: float,
        width: float,
        pad_top: float,
        pad_right: float,
        pad_bottom: float,
        pad_left: float,
        collected: List[Block],
    ):
        direction = block.style.get("flex-direction", "row")
        gap = parse_length(block.style.get("gap"), FLOW_GAP_X) or FLOW_GAP_X
        content_width = max(40.0, width - pad_left - pad_right)
        inner_left = left + pad_left
        inner_top = top + pad_top

        if direction == "column":
            cursor = inner_top
            max_width = 0
            start_index = len(collected)
            for child in block.children:
                cursor = self._layout_block(child, inner_left, cursor, content_width, collected)
                max_width = max(max_width, child.frame.width)
                cursor += gap
            height = (cursor - inner_top) + pad_bottom
            width_used = max_width + pad_left + pad_right
            block.frame = LayoutFrame(left, top, width or width_used, height)
            collected.insert(start_index, block)
        else:
            # row flex
            start_index = len(collected)
            count = max(1, len(block.children))
            explicit_sum = 0.0
            flex_count = 0
            for child in block.children:
                w = parse_length(child.style.get("width"), content_width)
                if w:
                    explicit_sum += w
                else:
                    flex_count += 1
            remaining = max(content_width - explicit_sum - (count - 1) * gap, 40.0)
            unit = remaining / max(1, flex_count) if flex_count else (remaining / count)
            cursor_x = inner_left
            max_height = 0
            for child in block.children:
                child_width = parse_length(child.style.get("width"), content_width) or unit
                self._layout_block(child, cursor_x, inner_top, child_width, collected)
                max_height = max(max_height, child.frame.height)
                cursor_x += child.frame.width + gap
            height = max_height + pad_top + pad_bottom
            width_used = cursor_x - inner_left + pad_right - gap if block.children else width
            block.frame = LayoutFrame(left, top, width or width_used, height)
            collected.insert(start_index, block)

    def _layout_grid(
        self,
        block: Block,
        left: float,
        top: float,
        width: float,
        pad_top: float,
        pad_right: float,
        pad_bottom: float,
        pad_left: float,
        collected: List[Block],
    ):
        template = block.style.get("grid-template-columns", "")
        parts = [p.strip() for p in template.split() if p.strip()]
        if not parts:
            parts = ["1fr", "1fr"]
        fr_units = sum(1 for p in parts if p.endswith("fr"))
        fixed_width = sum(parse_length(p, width) or 0 for p in parts if not p.endswith("fr"))
        remaining = max(40.0, width - pad_left - pad_right - fixed_width)
        fr_unit_width = remaining / max(1, fr_units) if fr_units else remaining / len(parts)
        col_widths = [(parse_length(p, width) if not p.endswith("fr") else fr_unit_width * float(p[:-2] or 1)) for p in parts]
        gap = parse_length(block.style.get("gap"), FLOW_GAP_X) or FLOW_GAP_X
        x = left + pad_left
        y = top + pad_top
        max_row_h = 0
        col_idx = 0
        start_index = len(collected)
        for child in block.children:
            cw = col_widths[col_idx % len(col_widths)]
            self._layout_block(child, x, y, cw, collected)
            max_row_h = max(max_row_h, child.frame.height)
            col_idx += 1
            if col_idx % len(col_widths) == 0:
                x = left + pad_left
                y += max_row_h + gap
                max_row_h = 0
            else:
                x += cw + gap
        if max_row_h:
            y += max_row_h
        height = (y - top) + pad_bottom
        block.frame = LayoutFrame(left, top, width, height)
        collected.insert(start_index, block)

    def _measure_block(self, block: Block, content_width: float) -> float:
        if block.kind == "text":
            _, h = self.measurer.measure_block(block.text, block.style, content_width)
            return h
        if block.kind == "list":
            font_size = parse_length(block.style.get("font-size")) or (DEFAULT_FONT_SIZE - 2)
            line_height = font_size * 1.3
            return len(block.items) * line_height + 12
        if block.kind == "table":
            font_size = parse_length(block.style.get("font-size")) or (DEFAULT_FONT_SIZE - 4)
            row_h = font_size * 1.5
            return len(block.table) * row_h + 20
        if block.kind in {"image", "chart"}:
            h = parse_length(block.style.get("height"), content_width * 0.6)
            if h is None:
                h = content_width * 0.6
            return h
        if block.kind == "group":
            return 0
        if block.kind == "shape":
            return parse_length(block.style.get("height"), 140.0) or 140.0
        return 120.0

    def _paginate(self, blocks: List[Block]) -> List[List[Block]]:
        slides: List[List[Block]] = []
        current: List[Block] = []
        offset = blocks[0].frame.top if blocks else 0.0
        for block in blocks:
            adjusted_top = block.frame.top - offset + BASE_MARGIN_Y
            bottom = adjusted_top + block.frame.height
            if bottom > self.height - BASE_MARGIN_Y and current:
                slides.append(current)
                current = []
                offset = block.frame.top
                adjusted_top = block.frame.top - offset + BASE_MARGIN_Y
            self._shift_block(block, adjusted_top - block.frame.top)
            current.append(block)
        if current:
            slides.append(current)
        return slides

    def _shift_block(self, block: Block, delta_y: float):
        block.frame.top += delta_y
        for child in block.children:
            self._shift_block(child, delta_y)


# --------------------------------------------------------------------------- #
# Chart parsing from <script> (Chart.js-like)
# --------------------------------------------------------------------------- #


class ChartBuilder:
    CHART_RE = re.compile(
        r"getElementById\(['\"](?P<cid>[^'\"]+)['\"]\).*?new\s+Chart\s*\([^,]+,\s*\{(?P<body>.*?)\}\s*\)",
        re.S,
    )
    LABEL_RE = re.compile(r"labels\s*:\s*\[(?P<labels>[^\]]+)\]", re.S)
    DATASET_RE = re.compile(r"datasets?\s*:\s*\[(?P<datasets>.*?)\]", re.S)
    SERIES_RE = re.compile(r"\{([^{}]*?)\}")
    TYPE_RE = re.compile(r"type\s*:\s*['\"](?P<type>[^'\"]+)['\"]", re.S)
    DATA_RE = re.compile(r"data\s*:\s*\[(?P<data>[^\]]+)\]", re.S)
    LABEL_FIELD_RE = re.compile(r"label\s*:\s*['\"](?P<label>[^'\"]+)['\"]")
    COLOR_RE = re.compile(r"color\s*:\s*['\"](?P<color>[^'\"]+)['\"]|backgroundColor\s*:\s*['\"](?P<bg>[^'\"]+)['\"]", re.S)

    def __init__(self, soup: BeautifulSoup):
        self.soup = soup

    def parse(self) -> List[ChartSpec]:
        specs: List[ChartSpec] = []
        for script in self.soup.find_all("script"):
            text = script.string or script.get_text()
            for match in self.CHART_RE.finditer(text):
                cid = match.group("cid")
                body = match.group("body")
                labels = self._parse_labels(body)
                series = self._parse_series(body)
                ctype = self._parse_type(body)
                if labels and series:
                    specs.append(ChartSpec(chart_type=ctype, labels=labels, series=series, target_ids=[cid]))
        return specs

    def _parse_labels(self, body: str) -> List[str]:
        m = self.LABEL_RE.search(body)
        if not m:
            return []
        raw = m.group("labels")
        return [self._strip_token(tok) for tok in raw.split(",") if self._strip_token(tok)]

    def _parse_series(self, body: str) -> List[ChartSeries]:
        series_list: List[ChartSeries] = []
        dm = self.DATASET_RE.search(body)
        if not dm:
            return series_list
        datasets_raw = dm.group("datasets")
        for sm in self.SERIES_RE.finditer(datasets_raw):
            snippet = sm.group(1)
            label = self._strip_token(self.LABEL_FIELD_RE.search(snippet).group("label")) if self.LABEL_FIELD_RE.search(snippet) else "Series"
            data_match = self.DATA_RE.search(snippet)
            if not data_match:
                continue
            values = [self._to_number(self._strip_token(tok)) for tok in data_match.group("data").split(",") if self._strip_token(tok)]
            color_match = self.COLOR_RE.search(snippet)
            color_val = None
            if color_match:
                color_val = color_match.group("color") or color_match.group("bg")
            series_list.append(ChartSeries(name=label, values=values, color=color_val))
        return series_list

    def _parse_type(self, body: str) -> str:
        m = self.TYPE_RE.search(body)
        if m:
            return m.group("type").strip().lower()
        return "bar"

    def _strip_token(self, token: str) -> str:
        return token.strip().strip("'\"")

    def _to_number(self, token: str) -> float:
        try:
            return float(token)
        except ValueError:
            return 0.0


# --------------------------------------------------------------------------- #
# PPTX rendering
# --------------------------------------------------------------------------- #


class PPTXRenderer:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_REF_WIDTH / PX_PER_INCH)
        self.prs.slide_height = Inches(SLIDE_REF_HEIGHT / PX_PER_INCH)
        self.blank = self.prs.slide_layouts[6]

    def render_slides(self, slide_blocks: List[List[Block]]):
        for blocks in slide_blocks:
            slide = self.prs.slides.add_slide(self.blank)
            for block in blocks:
                self._render_block(slide, block)

    def _render_block(self, slide, block: Block):
        if block.kind == "group":
            self._render_group(slide, block)
            return
        if block.kind == "text":
            self._render_text(slide, block)
        elif block.kind == "list":
            self._render_list(slide, block)
        elif block.kind == "table":
            self._render_table(slide, block)
        elif block.kind == "image":
            self._render_image(slide, block)
        elif block.kind == "chart":
            self._render_chart(slide, block)
        elif block.kind == "shape":
            self._render_shape(slide, block)

    def _geometry(self, frame: LayoutFrame) -> Tuple[float, float, float, float]:
        left = Inches(frame.left / PX_PER_INCH)
        top = Inches(frame.top / PX_PER_INCH)
        width = Inches(max(frame.width, 20.0) / PX_PER_INCH)
        height = Inches(max(frame.height, 20.0) / PX_PER_INCH)
        return left, top, width, height

    def _render_group(self, slide, block: Block):
        # background rectangle first if any
        if self._has_shape(block.style):
            left, top, width, height = self._geometry(block.frame)
            rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            self._apply_shape_style(rect, block.style)
        for child in block.children:
            self._render_block(slide, child)

    def _render_text(self, slide, block: Block):
        left, top, width, height = self._geometry(block.frame)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.clear()
        paragraph = tf.paragraphs[0]
        runs = block.runs or [TextRun(text=block.text)]
        current_paragraph = paragraph
        base_align = block.style.get("text-align", "").lower()
        for run_data in runs:
            pieces = run_data.text.split("\n")
            for idx, piece in enumerate(pieces):
                if idx > 0:
                    current_paragraph = tf.add_paragraph()
                run = current_paragraph.add_run()
                run.text = piece
                font = run.font
                font.size = pt_from_px(run_data.font_size or parse_length(block.style.get("font-size")) or DEFAULT_FONT_SIZE)
                font.bold = run_data.bold
                font.italic = run_data.italic
                color = parse_color(run_data.color) or parse_color(block.style.get("color"))
                if color:
                    font.color.rgb = color
        if base_align == "center":
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.CENTER
        elif base_align == "right":
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.RIGHT
        self._apply_shape_style(shape, block.style)

    def _render_list(self, slide, block: Block):
        left, top, width, height = self._geometry(block.frame)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.clear()
        for idx, item in enumerate(block.items):
            para = tf.add_paragraph() if idx else tf.paragraphs[0]
            para.text = item
            para.level = 0
            para.font.size = pt_from_px(parse_length(block.style.get("font-size")) or (DEFAULT_FONT_SIZE - 2))
            bullet = block.style.get("list-style-type", "disc")
            if bullet == "none":
                para.level = 1
        self._apply_shape_style(shape, block.style)

    def _render_table(self, slide, block: Block):
        if not block.table:
            return
        rows = len(block.table)
        cols = max(len(r) for r in block.table)
        left, top, width, height = self._geometry(block.frame)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        col_width_emu = int(width.emu / cols)
        row_height_emu = int(height.emu / rows)
        for c in range(cols):
            table.columns[c].width = col_width_emu
        for r in range(rows):
            table.rows[r].height = row_height_emu
        for r, row in enumerate(block.table):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = value
                style = {}
                if r < len(block.cell_styles) and c < len(block.cell_styles[r]):
                    style = block.cell_styles[r][c]
                para = cell.text_frame.paragraphs[0]
                para.font.size = pt_from_px(parse_length(style.get("font-size")) or parse_length(block.style.get("font-size")) or DEFAULT_FONT_SIZE - 4)
                para.font.bold = (style.get("font-weight") in {"bold", "700", "800"})
                color = parse_color(style.get("color")) or parse_color(block.style.get("color"))
                if color:
                    para.font.color.rgb = color
                align = style.get("text-align", "").lower()
                if align == "center":
                    para.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    para.alignment = PP_ALIGN.RIGHT
                bg = parse_color(style.get("background-color"))
                if bg:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = bg
                border = parse_color(style.get("border-color"))
                width_px = parse_length(style.get("border-width"), 1.2) or 1.2
                if border:
                    self._apply_cell_border(cell, border, width_px)

    def _apply_cell_border(self, cell, color: RGBColor, width_px: float):
        width_emu = pt_from_px(width_px).emu
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        border_tags = ("a:lnL", "a:lnR", "a:lnT", "a:lnB")
        for tag in border_tags:
            ln = tcPr.find(qn(tag))
            if ln is None:
                from pptx.oxml.xmlchemy import OxmlElement

                ln = OxmlElement(tag)
                tcPr.append(ln)
            ln.set("w", str(int(width_emu)))
            solid = ln.find(qn("a:solidFill"))
            if solid is None:
                from pptx.oxml.xmlchemy import OxmlElement

                solid = OxmlElement("a:solidFill")
                ln.append(solid)
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is None:
                from pptx.oxml.xmlchemy import OxmlElement

                srgb = OxmlElement("a:srgbClr")
                solid.append(srgb)
            srgb.set("val", "%02X%02X%02X" % (color[0], color[1], color[2]))

    def _render_image(self, slide, block: Block):
        left, top, width, height = self._geometry(block.frame)
        if block.image_path and block.image_path.exists():
            slide.shapes.add_picture(str(block.image_path), left, top, width=width, height=height)
        else:
            placeholder = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
            placeholder.text = block.text or "[image]"

    def _render_chart(self, slide, block: Block):
        left, top, width, height = self._geometry(block.frame)
        chart_data = self._chart_data_from_spec(block.chart)
        slide.shapes.add_chart(self._chart_type(block.chart), left, top, width, height, chart_data)

    def _chart_type(self, spec: Optional[ChartSpec]):
        if not spec:
            return XL_CHART_TYPE.COLUMN_CLUSTERED
        map_type = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }
        return map_type.get(spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def _chart_data_from_spec(self, spec: Optional[ChartSpec]) -> ChartData:
        data = CategoryChartData()
        if not spec:
            data.categories = ["A", "B", "C"]
            data.add_series("Series", (1, 2, 3))
            return data
        data.categories = spec.labels
        for series in spec.series:
            data.add_series(series.name, series.values)
        return data

    def _render_shape(self, slide, block: Block):
        left, top, width, height = self._geometry(block.frame)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
        self._apply_shape_style(shape, block.style)

    def _apply_shape_style(self, shape, style: Dict[str, str]):
        fill_color = parse_color(style.get("background") or style.get("background-color"))
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            if hasattr(shape, "fill"):
                shape.fill.background()
        border_color = parse_color(style.get("border-color") or self._parse_border_color(style.get("border")))
        border_width = (
            parse_length(style.get("border-width"))
            or self._parse_border_width(style.get("border"))
            or 2.0
        )
        if border_color and hasattr(shape, "line"):
            shape.line.color.rgb = border_color
            shape.line.width = pt_from_px(border_width)

    def _parse_border_color(self, border: Optional[str]) -> Optional[str]:
        if not border:
            return None
        parts = border.split()
        for p in parts:
            if p.startswith("#") or p.lower() in COLOR_KEYWORDS:
                return p
        return None

    def _parse_border_width(self, border: Optional[str]) -> Optional[float]:
        if not border:
            return None
        for part in border.split():
            val = parse_length(part)
            if val:
                return val
        return None

    def save(self, output: Path):
        self.prs.save(output)

    def _has_shape(self, style: Dict[str, str]) -> bool:
        return bool(style.get("background") or style.get("background-color") or style.get("border") or style.get("border-color"))


# --------------------------------------------------------------------------- #
# Conversion flow
# --------------------------------------------------------------------------- #


def resolve_slides(soup: BeautifulSoup) -> List[Tag]:
    slides = soup.select(".slide, .slide-container, [data-slide]")
    if slides:
        return slides
    if soup.body:
        return [soup.body]
    raise ValueError("No slide containers found. Use .slide or data-slide markers.")


def convert(input_html: Path, output_pptx: Path):
    soup = BeautifulSoup(input_html.read_text(encoding="utf-8"), "lxml")
    style_parser = StyleParser(soup)
    charts = ChartBuilder(soup).parse()
    builder = BlockBuilder(style_parser, input_html.parent, charts)
    measurer = TextMeasurer()
    layout_engine = LayoutEngine(SLIDE_REF_WIDTH, SLIDE_REF_HEIGHT, measurer)
    renderer = PPTXRenderer()

    for slide_tag in resolve_slides(soup):
        slide_style = style_parser.compute(slide_tag, None)
        # Respect slide-level width/height if provided.
        width = parse_length(slide_style.get("width"), SLIDE_REF_WIDTH) or SLIDE_REF_WIDTH
        height = parse_length(slide_style.get("height"), SLIDE_REF_HEIGHT) or SLIDE_REF_HEIGHT
        layout_engine.width = width
        layout_engine.height = height
        blocks = builder.build(slide_tag, slide_style)
        slide_blocks = layout_engine.layout_slide(blocks, slide_style)
        renderer.render_slides(slide_blocks)

    renderer.save(output_pptx)


def main():
    parser = argparse.ArgumentParser(description="Virtual-render HTML to PPTX without a browser.")
    parser.add_argument("input_html", type=Path)
    parser.add_argument("output_pptx", type=Path)
    args = parser.parse_args()
    convert(args.input_html, args.output_pptx)


if __name__ == "__main__":
    main()
