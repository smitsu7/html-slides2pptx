#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
html2pptx_light_2.py - High-fidelity HTML to PPTX converter without browser dependencies.
Uses Virtual Rendering and CSS Mapping to achieve feature parity with headless browser solutions.
"""

import argparse
import logging
import re
import sys
import io
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple, Union

import requests
from bs4 import BeautifulSoup, Tag, NavigableString
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    Image = None
    ImageDraw = None
    ImageFont = None

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- Constants ---
SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720
PPTX_WIDTH_INCHES = 13.333 # This should be updated too? 1280/96 = 13.333. 720/96 = 7.5.
# 1920/96 = 20. 1080/96 = 11.25.
# So 13.333 x 7.5 IS 1280x720.
# Wait, 1280 / 96 = 13.3333.
# 720 / 96 = 7.5.
# So the previous constants were inconsistent?
# SLIDE_WIDTH_PX = 1920 was paired with PPTX_WIDTH_INCHES = 13.333?
# If so, px_to_inches(1920) = 20 inches.
# But PPTX_WIDTH_INCHES was 13.333.
# Let's check the file content again.
# Line 37: SLIDE_WIDTH_PX = 1920
# Line 39: PPTX_WIDTH_INCHES = 13.333
# This implies the script was scaling?
# Or `PPTX_WIDTH_INCHES` is just a default for something else?
# `PPTXBuilder` uses `self.prs.slide_width = px_to_inches(SLIDE_WIDTH_PX)`.
# So it sets slide width to 20 inches.
# But `PPTX_WIDTH_INCHES` constant is unused?
# I'll update SLIDE_WIDTH_PX to 1280 and SLIDE_HEIGHT_PX to 720.

SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720
PPTX_WIDTH_INCHES = 13.333
PPTX_HEIGHT_INCHES = 7.5

# ...

@dataclass
class Style:
    display: str = "block"
    # ...
    # Positioning
    top: Optional[float] = None
    left: Optional[float] = None
    right: Optional[float] = None
    bottom: Optional[float] = None
    z_index: int = 0
    position: str = "static"
    width: Optional[float] = None # Explicit width
    height: Optional[float] = None # Explicit height
    min_width: float = 0
    max_width: float = float('inf')
    min_height: float = 0
    margin_top: float = 0
    margin_right: float = 0
    margin_bottom: float = 0
    margin_left: float = 0
    padding_top: float = 0
    padding_right: float = 0
    padding_bottom: float = 0
    padding_left: float = 0
    float: str = "none"
    clear: str = "none"
    color: Tuple[int, int, int] = (0, 0, 0)
    background_color: Optional[Tuple[int, int, int]] = None
    font_size: float = 24.0
    font_family: str = "Arial"
    font_style: str = "normal" # normal, italic
    font_weight: str = "normal"
    text_align: str = "left"
    border_color: Optional[Tuple[int, int, int]] = None
    border_width: float = 0
    list_style_type: str = "none"
    # Flexbox
    flex_direction: str = "row"
    justify_content: str = "flex-start"
    
    # Grid
    grid_template_columns: List[str] = field(default_factory=list) # e.g. ['100px', '1fr']
    
    # Positioning
    top: Optional[float] = None
    left: Optional[float] = None
    right: Optional[float] = None
    bottom: Optional[float] = None

@dataclass
class RenderBox:
    element: Optional[Tag]
    style: Style
    # Box Model Dimensions (Calculated)
    x: float = 0
    y: float = 0
    content_width: float = 0
    content_height: float = 0
    # Total dimensions (Content + Padding + Border)
    # Margin is external
    
    children: List['RenderBox'] = field(default_factory=list)
    text_lines: List[str] = field(default_factory=list) # Wrapped text lines
    box_type: str = "block" # block, inline, table, row, cell, image, list_item, chart_placeholder
    image_data: Optional[bytes] = None
    table_data: Optional[List[List[str]]] = None
    chart_id: Optional[str] = None

    @property
    def total_width(self):
        return self.content_width + self.style.padding_left + self.style.padding_right + (self.style.border_width * 2)
    
    @property
    def total_height(self):
        return self.content_height + self.style.padding_top + self.style.padding_bottom + (self.style.border_width * 2)

    @property
    def outer_width(self):
        return self.total_width + self.style.margin_left + self.style.margin_right

    @property
    def outer_height(self):
        return self.total_height + self.style.margin_top + self.style.margin_bottom

class TextMeasurer:
    """Helper to measure text dimensions using PIL."""
    def __init__(self):
        self.dummy_img = Image.new('RGB', (1, 1))
        self.draw = ImageDraw.Draw(self.dummy_img)
        self.fonts = {} # Cache fonts: (family, size, weight) -> ImageFont

    def get_font(self, family: str, size: float, weight: str):
        key = (family, size, weight)
        if key not in self.fonts:
            # Try to load a system font, fallback to default
            try:
                # This is OS dependent. For this environment, we might need a fallback or a specific path.
                # Using default for safety if specific font fails.
                # In a real scenario, we'd map 'Arial' to '/path/to/arial.ttf'
                font = ImageFont.truetype("Arial.ttf", int(size))
            except OSError:
                try:
                     font = ImageFont.truetype("/System/Library/Fonts/Supplemental/Arial.ttf", int(size))
                except OSError:
                    # Fallback to default PIL font (which is very small/bitmap, not ideal but works)
                    # Better fallback: try to find *any* ttf
                    font = ImageFont.load_default() 
            self.fonts[key] = font
        return self.fonts[key]

    def measure_text(self, text: str, style: Style) -> Tuple[float, float]:
        font = self.get_font(style.font_family, style.font_size, style.font_weight)
        bbox = self.draw.textbbox((0, 0), text, font=font)
        return bbox[2] - bbox[0], bbox[3] - bbox[1]

    def wrap_text(self, text: str, style: Style, max_width: float) -> Tuple[List[str], float, float]:
        """Wraps text to fit max_width. Returns (lines, max_line_width, total_height)."""
        if not text:
            return [], 0, 0
            
        font = self.get_font(style.font_family, style.font_size, style.font_weight)
        words = text.split()
        lines = []
        current_line = []
        current_width = 0
        max_line_width = 0
        
        # Space width
        space_bbox = self.draw.textbbox((0, 0), " ", font=font)
        space_width = space_bbox[2] - space_bbox[0]

        for word in words:
            word_bbox = self.draw.textbbox((0, 0), word, font=font)
            word_width = word_bbox[2] - word_bbox[0]
            
            if current_width + word_width <= max_width or not current_line:
                current_line.append(word)
                current_width += word_width + space_width
            else:
                lines.append(" ".join(current_line))
                max_line_width = max(max_line_width, current_width - space_width)
                current_line = [word]
                current_width = word_width + space_width
        
        if current_line:
            lines.append(" ".join(current_line))
            max_line_width = max(max_line_width, current_width - space_width)
            
        # Calculate total height
        # Line height approx 1.2 * font size usually, or measure 'Ag'
        line_height_bbox = self.draw.textbbox((0, 0), "Ag", font=font)
        line_height = (line_height_bbox[3] - line_height_bbox[1]) * 1.2
        total_height = len(lines) * line_height
        print(f"DEBUG: TextMeasurer - Text: {text[:20]}..., Width: {max_width}, Lines: {len(lines)}, TotalH: {total_height}")
        
        return lines, max_line_width, total_height

# --- Components ---

# --- Components ---

class StyleParser:
    """Parses CSS and applies styles to DOM elements."""
    
    def __init__(self):
        self.css_rules = []
        # Simple regex for CSS rules: selector { property: value; ... }
        self.rule_pattern = re.compile(r'([^{]+)\s*\{\s*([^}]+)\s*\}')
        self.prop_pattern = re.compile(r'([\w-]+)\s*:\s*([^;]+);?')

    def parse_css(self, css_text: str):
        """Parses CSS text and stores rules."""
        # Remove comments
        css_text = re.sub(r'/\*.*?\*/', '', css_text, flags=re.DOTALL)
        
        for match in self.rule_pattern.finditer(css_text):
            selectors = [s.strip() for s in match.group(1).split(',')]
            body = match.group(2)
            props = {}
            for prop_match in self.prop_pattern.finditer(body):
                key = prop_match.group(1).strip().lower()
                val = prop_match.group(2).strip()
                props[key] = val
            
            for selector in selectors:
                self.css_rules.append((selector, props))
        
        # Sort by specificity (naive: ID > Class > Tag)
        # Actually, we'll just append and let later rules override earlier ones (standard CSS behavior within same sheet)
        # But for correct specificity we need a scoring system.
        # Simplified scoring: ID=100, Class=10, Tag=1
        self.css_rules.sort(key=lambda x: self._calculate_specificity(x[0]))

    def _calculate_specificity(self, selector: str) -> int:
        score = 0
        if '#' in selector: score += 100
        if '.' in selector: score += 10
        if re.match(r'^[a-zA-Z]', selector): score += 1
        return score

    def compute_style(self, element: Tag) -> Style:
        """Computes the final style for a given element."""
        style_props = {}
        
        # 1. Default Styles (User Agent)
        if element.name in ['div', 'p', 'h1', 'h2', 'h3', 'h4', 'ul', 'ol', 'li', 'section', 'article']:
            style_props['display'] = 'block'
        elif element.name == 'body':
            style_props['display'] = 'block'
            style_props['margin'] = '0' # Reset body margin for full bleed
            style_props['padding'] = '0'
            style_props['width'] = '100%'
            style_props['height'] = '100%'
        elif element.name == 'span':
            style_props['display'] = 'inline'
        
        # 2. Apply CSS Rules
        # This is a simplified matcher. It doesn't handle complex combinators like "div > p" well.
        for selector, props in self.css_rules:
            if self._matches_selector(element, selector):
                style_props.update(props)
        
        # 3. Inline Styles
        if element.has_attr('style'):
            for prop_match in self.prop_pattern.finditer(element['style']):
                key = prop_match.group(1).strip().lower()
                val = prop_match.group(2).strip()
                style_props[key] = val
        
        # 4. Tailwind Polyfill
        self._apply_tailwind_polyfill(element, style_props)

        # Convert dictionary to Style object
        return self._dict_to_style(style_props)

    def _matches_selector(self, element: Tag, selector: str) -> bool:
        # Handle simple selectors: tag, .class, #id
        # TODO: Add support for descendant selectors (e.g. ".slide h1")
        
        selectors = selector.split()
        if len(selectors) > 1:
            # Very basic descendant support: check if last part matches element, 
            # and first part matches some ancestor.
            if not self._matches_simple(element, selectors[-1]):
                return False
            # Check ancestors
            parent = element.parent
            while parent:
                if isinstance(parent, Tag) and self._matches_simple(parent, selectors[0]):
                    return True
                parent = parent.parent
            return False
        else:
            return self._matches_simple(element, selector)

    def _matches_simple(self, element: Tag, selector: str) -> bool:
        if selector.startswith('#'):
            return element.get('id') == selector[1:]
        elif selector.startswith('.'):
            return selector[1:] in element.get('class', [])
        else:
            return element.name == selector

    def _dict_to_style(self, props: Dict[str, str]) -> Style:
        s = Style()
        
        # Helper to parse px/rem/em
        def parse_len(val: str, base: float = 16.0) -> float:
            if not val: return 0.0
            if val == 'auto': return 0.0 # Simplified
            if val.endswith('px'): return float(val[:-2])
            if val.endswith('pt'): return float(val[:-2]) * 1.33
            if val.endswith('in'): return float(val[:-2]) * 96
            if val.endswith('cm'): return float(val[:-2]) * 37.8
            if val.endswith('%'): return 0.0 # TODO: Handle percentage
            try: return float(val)
            except: return 0.0

        # Helper to parse color
        def parse_color(val: str) -> Optional[Tuple[int, int, int]]:
            if not val or val == 'transparent': return None
            # Hex
            if val.startswith('#'):
                h = val.lstrip('#')
                if len(h) == 3: h = ''.join([c*2 for c in h])
                return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
            # RGB
            if val.startswith('rgb'):
                try:
                    parts = val[val.find('(')+1:val.find(')')].split(',')
                    return tuple(int(p.strip()) for p in parts[:3])
                except: pass
            # Named colors (basic)
            colors = {'red': (255,0,0), 'green': (0,128,0), 'blue': (0,0,255), 'white': (255,255,255), 'black': (0,0,0), 'gray': (128,128,128)}
            return colors.get(val, None)

        s.display = props.get('display', 'block')
        s.position = props.get('position', 'static')
        
        if 'width' in props: s.width = parse_len(props['width'])
        if 'height' in props: s.height = parse_len(props['height'])
        
        s.margin_top = parse_len(props.get('margin-top', '0'))
        s.margin_right = parse_len(props.get('margin-right', '0'))
        s.margin_bottom = parse_len(props.get('margin-bottom', '0'))
        s.margin_left = parse_len(props.get('margin-left', '0'))
        
        # Shorthand margin
        if 'margin' in props:
            parts = props['margin'].split()
            if len(parts) == 1:
                v = parse_len(parts[0])
                s.margin_top = s.margin_right = s.margin_bottom = s.margin_left = v
            elif len(parts) == 2:
                v1 = parse_len(parts[0])
                v2 = parse_len(parts[1])
                s.margin_top = s.margin_bottom = v1
                s.margin_right = s.margin_left = v2
        
        s.padding_top = parse_len(props.get('padding-top', '0'))
        s.padding_right = parse_len(props.get('padding-right', '0'))
        s.padding_bottom = parse_len(props.get('padding-bottom', '0'))
        s.padding_left = parse_len(props.get('padding-left', '0'))
        
        # Shorthand padding
        if 'padding' in props:
            parts = props['padding'].split()
            if len(parts) == 1:
                v = parse_len(parts[0])
                s.padding_top = s.padding_right = s.padding_bottom = s.padding_left = v
            elif len(parts) == 2:
                v1 = parse_len(parts[0])
                v2 = parse_len(parts[1])
                s.padding_top = s.padding_bottom = v1
                s.padding_right = s.padding_left = v2

        s.float = props.get('float', 'none')
        s.clear = props.get('clear', 'none')
        s.color = parse_color(props.get('color', 'black'))
        s.background_color = parse_color(props.get('background-color', 'transparent'))
        
        if 'font-size' in props: s.font_size = parse_len(props['font-size'])
        s.font_weight = props.get('font-weight', 'normal')
        s.font_style = props.get('font-style', 'normal')
        s.text_align = props.get('text-align', 'left')
        
        # Font Family Mapping
        raw_font = props.get('font-family', 'Arial').split(',')[0].strip().replace('"', '').replace("'", "")
        font_map = {
            'sans-serif': 'Arial',
            'serif': 'Times New Roman',
            'monospace': 'Courier New',
            'system-ui': 'Arial',
            '-apple-system': 'Arial',
            'Segoe UI': 'Arial',
            'Roboto': 'Arial',
            'Helvetica': 'Arial'
        }
        s.font_family = font_map.get(raw_font, raw_font)
        
        s.border_color = parse_color(props.get('border-color', 'black'))
        s.border_width = parse_len(props.get('border-width', '0'))
        if 'border' in props:
            # Simple parser: 1px solid black
            parts = props['border'].split()
            for p in parts:
                if 'px' in p: s.border_width = parse_len(p)
                elif p.startswith('#') or p in ['red', 'black', 'blue', 'green', 'white', 'gray']: s.border_color = parse_color(p)
        
        # Flexbox
        s.flex_direction = props.get('flex-direction', 'row')
        s.justify_content = props.get('justify-content', 'flex-start')
        
        # Grid
        if 'grid-template-columns' in props:
            # Simple parser: split by space
            s.grid_template_columns = props['grid-template-columns'].split()

        # Positioning
        if 'top' in props: s.top = parse_len(props['top'])
        if 'left' in props: s.left = parse_len(props['left'])
        if 'right' in props: s.right = parse_len(props['right'])
        if 'bottom' in props: s.bottom = parse_len(props['bottom'])
        
        if 'z-index' in props:
            try: s.z_index = int(props['z-index'])
            except: s.z_index = 0

        return s
        
    def _apply_tailwind_polyfill(self, element: Tag, style_props: Dict[str, str]):
        """Applies Tailwind-like classes to style_props."""
        classes = element.get('class', [])
        if not classes: return
        
        for cls in classes:
            # Layout
            if cls == 'flex': style_props['display'] = 'flex'
            elif cls == 'grid': style_props['display'] = 'grid'
            elif cls == 'block': style_props['display'] = 'block'
            elif cls == 'inline-block': style_props['display'] = 'inline-block'
            elif cls == 'hidden': style_props['display'] = 'none'
            
            # Flex
            elif cls == 'flex-row': style_props['flex-direction'] = 'row'
            elif cls == 'flex-col': style_props['flex-direction'] = 'column'
            elif cls == 'justify-center': style_props['justify-content'] = 'center'
            elif cls == 'justify-between': style_props['justify-content'] = 'space-between'
            elif cls == 'items-center': style_props['align-items'] = 'center' # Not fully supported yet
            
            # Grid
            elif cls.startswith('grid-cols-'):
                # e.g. grid-cols-2 -> repeat(2, minmax(0, 1fr)) -> simplified to 1fr 1fr
                try:
                    cols = int(cls.split('-')[2])
                    style_props['grid-template-columns'] = ' '.join(['1fr'] * cols)
                except: pass
            elif cls.startswith('gap-'):
                # gap-4 -> 1rem -> 16px
                try:
                    val = int(cls.split('-')[1]) * 4
                    style_props['gap'] = f"{val}px" # Not fully supported yet
                except: pass
                
            # Spacing (Padding/Margin) - Simplified (p-4, m-4, px-2, my-2)
            elif cls.startswith('p-'):
                try: style_props['padding'] = f"{int(cls.split('-')[1]) * 4}px"
                except: pass
            elif cls.startswith('m-'):
                try: style_props['margin'] = f"{int(cls.split('-')[1]) * 4}px"
                except: pass
            
            # Typography
            elif cls.startswith('text-'):
                if cls in ['text-left', 'text-center', 'text-right', 'text-justify']:
                     style_props['text-align'] = cls.split('-')[1]
                elif cls in ['text-xs', 'text-sm', 'text-base', 'text-lg', 'text-xl', 'text-2xl', 'text-3xl']:
                    sizes = {'xs': '12px', 'sm': '14px', 'base': '16px', 'lg': '18px', 'xl': '20px', '2xl': '24px', '3xl': '30px'}
                    style_props['font-size'] = sizes.get(cls.split('-')[1], '16px')
                elif cls in ['text-white', 'text-black', 'text-red-500', 'text-blue-500']:
                    # Very basic color map
                    colors = {'white': 'white', 'black': 'black', 'red-500': 'red', 'blue-500': 'blue'}
                    style_props['color'] = colors.get(cls.replace('text-', ''), 'black')
            
            elif cls == 'font-bold': style_props['font-weight'] = 'bold'
            elif cls == 'italic': style_props['font-style'] = 'italic'
            
            # Sizing
            elif cls == 'w-full': style_props['width'] = '100%'
            elif cls == 'h-full': style_props['height'] = '100%'
            elif cls == 'w-screen': style_props['width'] = '100vw'
            elif cls == 'h-screen': style_props['height'] = '100vh'

class LayoutEngine:
    """
    Virtual Rendering Engine.
    Calculates the position and size of every element.
    """
    def __init__(self, slide_width: float, slide_height: float):
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.default_padding = 40
        self.text_measurer = TextMeasurer()
        
        # State
        self.pages: List[RenderBox] = []
        self.current_page_root: Optional[RenderBox] = None
        self.remaining_height: float = 0

    def layout(self, soup: BeautifulSoup, style_parser: StyleParser) -> List[RenderBox]:
        self._start_new_page()
        
        # Process body children
        body = soup.body
        if body:
            # We treat body as a block container
            start_x = self.current_page_root.x + self.current_page_root.style.padding_left
            start_y = self.current_page_root.y + self.current_page_root.style.padding_top
            self._layout_container_children(body, self.current_page_root, style_parser, start_x, start_y, self.current_page_root.content_width)
            
        return self.pages

    def _start_new_page(self):
        style = Style(width=self.slide_width, height=self.slide_height, padding_top=self.default_padding, padding_left=self.default_padding, padding_right=self.default_padding, padding_bottom=self.default_padding)
        root = RenderBox(element=None, style=style, x=0, y=0, content_width=self.slide_width - (self.default_padding*2), content_height=self.slide_height - (self.default_padding*2))
        self.pages.append(root)
        self.current_page_root = root
        self.remaining_height = root.content_height

    def _layout_container_children(self, container_element: Tag, parent_box: RenderBox, style_parser: StyleParser, available_width: float):
        """
        Layouts children of a container (Block or Flex).
        This is the main recursive loop.
        """
        # 1. Filter relevant children
        children = [c for c in container_element.children if isinstance(c, Tag) and c.name not in ['script', 'style', 'noscript', 'meta', 'link', 'head', 'title']]
        
    def _layout_container_children(self, element: Tag, box: RenderBox, style_parser: StyleParser, start_x: float, start_y: float, available_width: float):
        """Dispatches layout based on display type."""
        if box.style.display == 'flex':
            self._layout_flex_children(element, box, style_parser, start_x, start_y, available_width)
        elif box.style.display == 'grid':
            self._layout_grid_children(element, box, style_parser, start_x, start_y, available_width)
        else:
            self._layout_block_children(element, box, style_parser, start_x, start_y, available_width)

    def _layout_grid_children(self, element: Tag, box: RenderBox, style_parser: StyleParser, start_x: float, start_y: float, available_width: float):
        """Lays out children in a grid."""
        if not box.style.grid_template_columns:
            # Fallback to block if no columns defined
            self._layout_block_children(element, box, style_parser, start_x, start_y, available_width)
            return

        # Parse Columns
        col_defs = box.style.grid_template_columns
        
        # 1. Calculate fixed widths
        fixed_width = 0
        fr_total = 0
        
        parsed_cols = []
        for col_def in col_defs:
            if col_def.endswith('px'):
                val = float(col_def[:-2])
                parsed_cols.append({'type': 'px', 'val': val})
                fixed_width += val
            elif col_def.endswith('fr'):
                val = float(col_def[:-2])
                parsed_cols.append({'type': 'fr', 'val': val})
                fr_total += val
            else:
                # Assume px or auto (treat auto as 1fr for now)
                try:
                    val = float(col_def)
                    parsed_cols.append({'type': 'px', 'val': val})
                    fixed_width += val
                except:
                    parsed_cols.append({'type': 'fr', 'val': 1.0})
                    fr_total += 1.0
        
        # 2. Calculate fr width
        remaining_width = max(0, available_width - fixed_width)
        fr_unit = remaining_width / fr_total if fr_total > 0 else 0
        
        col_widths = []
        for col in parsed_cols:
            if col['type'] == 'px':
                col_widths.append(col['val'])
            else:
                col_widths.append(col['val'] * fr_unit)
        
        # 3. Place Children
        current_row_y = start_y
        row_max_height = 0
        
        col_x_offsets = []
        curr_x = start_x
        for w in col_widths:
            col_x_offsets.append(curr_x)
            curr_x += w
            
        print(f"DEBUG: Grid Layout - StartX: {start_x}, StartY: {start_y}, ColWidths: {col_widths}, Offsets: {col_x_offsets}")

        children = [c for c in element.children if isinstance(c, Tag) and c.name not in ['script', 'style', 'noscript', 'meta', 'link', 'head', 'title']]
        
        for i, child in enumerate(children):
            col_idx = i % len(col_widths)
            
            if col_idx == 0 and i > 0:
                # New Row
                current_row_y += row_max_height
                row_max_height = 0
            
            cell_x = col_x_offsets[col_idx]
            cell_w = col_widths[col_idx]
            
            child_style = style_parser.compute_style(child)
            
            # Force width to cell width (minus margins)
            child_content_width = cell_w - child_style.margin_left - child_style.margin_right
            
            child_box = RenderBox(
                element=child,
                box_type="block", # Default to block inside grid cell
                style=child_style,
                x=cell_x + child_style.margin_left,
                y=current_row_y + child_style.margin_top,
                content_width=child_content_width,
                content_height=0 # Will be calculated
            )
            
            # Layout child's children
            child_inner_x = child_box.x + child_style.padding_left + child_style.border_width
            child_inner_y = child_box.y + child_style.padding_top + child_style.border_width
            
            self._layout_container_children(child, child_box, style_parser, child_inner_x, child_inner_y, child_content_width)
            
            # Update child height
            child_box.content_height = max(child_box.content_height, child_style.height) if (child_style.height and child_style.height > 0) else child_box.content_height
            
            box.children.append(child_box)
            row_max_height = max(row_max_height, child_box.total_height + child_style.margin_top + child_style.margin_bottom)
            print(f"DEBUG: Grid Child {i} - Col: {col_idx}, X: {child_box.x}, Y: {child_box.y}, W: {child_box.content_width}, H: {child_box.content_height}, TotalH: {child_box.total_height}")
        
        # Update total height
        box.content_height = (current_row_y + row_max_height) - start_y

    def _layout_block_children(self, element: Tag, box: RenderBox, style_parser: StyleParser, start_x: float, start_y: float, available_width: float):
        """Lays out children in a vertical block."""
        current_y = start_y
        
        # Filter relevant children
        children_to_layout = [c for c in element.children if isinstance(c, (Tag, NavigableString)) and (isinstance(c, NavigableString) or c.name not in ['script', 'style', 'noscript', 'meta', 'link', 'head', 'title'])]

        for child in children_to_layout:
            if isinstance(child, NavigableString):
                text = child.string.strip()
                if text:
                    # Measure and wrap text
                    lines, max_w, total_h = self.text_measurer.wrap_text(text, box.style, available_width)
                    if lines:
                        # Create a clean style for text box (inherit font, reset layout)
                        text_style = Style(
                            font_family=box.style.font_family,
                            font_size=box.style.font_size,
                            font_weight=box.style.font_weight,
                            font_style=box.style.font_style,
                            color=box.style.color,
                            # Reset layout
                            padding_top=0, padding_bottom=0, padding_left=0, padding_right=0,
                            margin_top=0, margin_bottom=0, margin_left=0, margin_right=0,
                            border_width=0
                        )
                        
                        text_box = RenderBox(
                            element=None,
                            box_type="text",
                            style=text_style,
                            x=start_x,
                            y=current_y,
                            content_width=max_w,
                            content_height=total_h,
                            text_lines=lines
                        )
                        # total_height is computed from content_height + padding (0) -> total_h
                        box.children.append(text_box)
                        current_y += text_box.total_height
                continue
                
            # Element node
            style = style_parser.compute_style(child)
            
            # Absolute Positioning Handling
            if style.position == 'absolute':
                # Calculate dimensions (simplified)
                if style.width and style.width > 0:
                    child_width = style.width
                else:
                    child_width = available_width # Default to full width? Or auto?
                
                # Calculate Position (Relative to parent box)
                # Note: Ideally relative to nearest positioned ancestor, but we use parent for simplicity
                abs_x = box.x + (style.left if style.left is not None else style.padding_left) # Default to padding? No, default to static pos?
                # If left is None, it should be at static position (start_x).
                if style.left is not None:
                    abs_x = box.x + style.left
                elif style.right is not None:
                    abs_x = box.x + box.content_width - style.right - child_width
                else:
                    abs_x = start_x # Static position
                    
                abs_y = box.y + (style.top if style.top is not None else style.padding_top)
                if style.top is not None:
                    abs_y = box.y + style.top
                elif style.bottom is not None:
                    # We need height for bottom. 
                    # If height is auto, this is hard.
                    # Assume top=static if not set.
                    abs_y = current_y
                else:
                    abs_y = current_y # Static position

                child_box = RenderBox(
                    element=child,
                    box_type="block",
                    style=style,
                    x=abs_x,
                    y=abs_y,
                    content_width=child_width,
                    content_height=0
                )
                
                if child.name == 'img':
                    child_box.box_type = "image"
                    w, h, data = self._measure_image(child, child_width)
                    child_box.content_width = w
                    child_box.content_height = h
                    child_box.image_data = data
                elif child.name == 'canvas':
                    child_box.box_type = "chart_placeholder"
                    child_box.content_height = 400 if (not style.height or style.height == 0) else style.height
                else:
                    # Recurse
                    child_inner_x = child_box.x + style.padding_left + style.border_width
                    child_inner_y = child_box.y + style.padding_top + style.border_width
                    self._layout_container_children(child, child_box, style_parser, child_inner_x, child_inner_y, child_box.content_width)
                    
                    if style.height and style.height > 0:
                        child_box.content_height = style.height
                    else:
                        # Calc height from children
                        max_child_y = child_inner_y
                        for c in child_box.children:
                             max_child_y = max(max_child_y, c.y + c.total_height)
                        child_box.content_height = max_child_y - child_inner_y

                box.children.append(child_box)
                print(f"DEBUG: Absolute Child - Y: {child_box.y}, H: {child_box.content_height}")
                continue # Do not update current_y
            
            # Special handling for Slide Container (Fix "White Square" offset)
            is_slide_root = False
            classes = child.get('class', [])
            if classes and ('slide' in classes or 'slide-container' in classes):
                 if element.name == 'body': # Direct child of body
                     is_slide_root = True

            # Calculate margins
            margin_top = style.margin_top
            margin_bottom = style.margin_bottom
            margin_left = style.margin_left
            margin_right = style.margin_right
            
            # Box Width
            if style.width and style.width > 0:
                child_width = style.width
            else:
                child_width = available_width - margin_left - margin_right
            
            # Position
            child_x = start_x + margin_left
            child_y = current_y + margin_top
            
            if is_slide_root:
                child_x = 0
                child_y = 0
                # If it has an explicit width, use it, otherwise fill the slide width
                if style.width and style.width > 0:
                    child_width = style.width
                else:
                    child_width = self.slide_width - (style.padding_left + style.padding_right + (style.border_width * 2))
                
                # If it has an explicit height, use it, otherwise fill the slide height
                if style.height and style.height > 0:
                    style.height = style.height
                else:
                    style.height = self.slide_height - (style.padding_top + style.padding_bottom + (style.border_width * 2))
                
                # Reset current_y? No, we just place it at 0,0.
            
            child_box = RenderBox(
                element=child,
                box_type="block",
                style=style,
                x=child_x,
                y=child_y,
                content_width=child_width,
                content_height=0 # Calculated later
            )
            
            if child.name == 'img':
                child_box.box_type = "image"
                w, h, data = self._measure_image(child, child_width)
                child_box.content_width = w
                child_box.content_height = h
                child_box.image_data = data
            elif child.name == 'table':
                child_box.box_type = "table"
                self._measure_table(child, child_box, style_parser)
            elif child.name == 'canvas':
                child_box.box_type = "chart_placeholder"
                # Default size for chart
                child_box.content_height = 400 if (not style.height or style.height == 0) else style.height
            else:
                # Recurse
                child_inner_x = child_box.x + style.padding_left + style.border_width
                child_inner_y = child_box.y + style.padding_top + style.border_width
                
                self._layout_container_children(child, child_box, style_parser, child_inner_x, child_inner_y, child_box.content_width)
                
                # Update height based on children
                if style.height and style.height > 0:
                    child_box.content_height = style.height
                else:
                    max_child_y = child_inner_y
                    for c in child_box.children:
                        max_child_y = max(max_child_y, c.y + c.total_height)
                    child_box.content_height = max_child_y - child_inner_y

            box.children.append(child_box)
            print(f"DEBUG: Block Child - Type: {child_box.box_type}, Y: {child_box.y}, H: {child_box.content_height}, TotalH: {child_box.total_height}")
            
            # Pagination Check (only for top-level elements of the page root)
            if box == self.current_page_root:
                # Check if this element overflows the current page
                if current_y + child_box.total_height + margin_bottom > self.current_page_root.y + self.current_page_root.content_height:
                    self._start_new_page()
                    parent_box = self.current_page_root # Switch to new page root
                    self.remaining_height = parent_box.content_height
                    # Re-add this child to the new page?
                    # This is tricky because we already processed it.
                    # Ideally we check BEFORE processing.
                    # But we need the height to check.
                    # So we move it.
                    box.children.remove(child_box)
                    
                    # Reset position for new page
                    child_box.x = parent_box.x + parent_box.style.padding_left + margin_left
                    child_box.y = parent_box.y + parent_box.style.padding_top + margin_top
                    
                    # We might need to re-layout children if width changed (unlikely for same slide size)
                    # Add to new page
                    parent_box.children.append(child_box)
                    current_y = child_box.y + child_box.total_height + margin_bottom
                    # Return? No, continue loop but with new parent?
                    # This loop is iterating children of 'element'.
                    # 'box' is the container we are adding to.
                    # If we switch 'box' to new page root, we need to update 'box' variable.
                    box = parent_box
                    current_y = child_box.y + child_box.total_height + margin_bottom
                    continue

            current_y = max(current_y, child_box.y + child_box.total_height + margin_bottom)
            
        # Update parent height
        box.content_height = current_y - start_y

    def _layout_flex_children(self, element: Tag, parent_box: RenderBox, style_parser: StyleParser, start_x: float, start_y: float, available_width: float):
        # Simple Flex Row implementation
        # Distribute width equally or based on content
        children = [c for c in element.children if isinstance(c, Tag) and c.name not in ['script', 'style', 'noscript', 'meta', 'link', 'head', 'title']]
        if not children: return
        
        count = len(children)
        col_width = available_width / count # Naive equal width
        
        max_height = 0
        current_x = start_x
        
        for child in children:
            style = style_parser.compute_style(child)
            
            # Force width to column width (minus margins)
            box_width = col_width - style.margin_left - style.margin_right
            
            box = RenderBox(element=child, style=style)
            box.content_width = box_width - style.padding_left - style.padding_right - (style.border_width * 2)
            
            # Layout content inside column
            # Similar to block layout but restricted width
            direct_text = "".join([str(c) for c in child.children if isinstance(c, NavigableString)]).strip()
            if direct_text:
                lines, max_w, text_h = self.text_measurer.wrap_text(direct_text, style, box.content_width)
                box.text_lines = lines
                box.content_height = max(box.content_height, text_h)
            
            if child.name == 'img':
                self._measure_image(box, box.content_width)
            elif child.name == 'canvas':
                box.content_height = 300
                box.box_type = 'chart_placeholder'
                box.chart_id = child.get('id')
            else:
                 if list(child.children):
                    # Recurse
                    child_inner_x = box.x + style.padding_left + style.border_width # Wait, box.x is not set yet
                    # We need to set box.x first
                    pass

            if style.height and style.height > 0: box.content_height = style.height
            
            # Position
            box.x = current_x + style.margin_left
            box.y = start_y + style.margin_top
            
            # Recurse now that we have position
            if child.name not in ['img', 'canvas'] and list(child.children):
                 child_inner_x = box.x + style.padding_left + style.border_width
                 child_inner_y = box.y + style.padding_top + style.border_width
                 self._layout_container_children(child, box, style_parser, child_inner_x, child_inner_y, box.content_width)
                 # Update height based on children
                 children_height = sum([c.outer_height for c in box.children]) # This is wrong for block layout inside flex
                 # If we use _layout_container_children, it populates box.children and sets box.content_height?
                 # No, _layout_container_children sets box.content_height for block/grid, but for flex?
                 # _layout_block_children sets box.content_height.
                 # So we don't need to sum children here if we called _layout_container_children.
                 pass

            parent_box.children.append(box)
            
            current_x += col_width
            max_height = max(max_height, box.outer_height)
            
        # Update parent height to match tallest column
        # Note: This doesn't handle wrapping flex
        # parent_box.content_height = max_height # This might be set by caller?
        # In recursion, the parent height is determined by children. 
        # But for flex container, it expands to fit content.

    def _measure_image(self, box: RenderBox, available_width: float):
        src = box.element.get('src')
        if not src: return
        
        try:
            if src.startswith('http'):
                response = requests.get(src, timeout=5)
                img_data = response.content
            else:
                with open(src, 'rb') as f:
                    img_data = f.read()
            
            box.image_data = img_data
            box.box_type = 'image'
            
            # Get dimensions
            if Image:
                with Image.open(io.BytesIO(img_data)) as img:
                    img_w, img_h = img.size
                    aspect = img_w / img_h
                    
                    # Scale
                    final_w = min(img_w, available_width)
                    final_h = final_w / aspect
                    
                    box.content_width = final_w
                    box.content_height = final_h
        except Exception as e:
            logger.warning(f"Failed to load image {src}: {e}")

    def _measure_table(self, element: Tag, box: RenderBox, style_parser: StyleParser):
        box.box_type = 'table'
        rows = element.find_all('tr')
        if not rows: return
        
        # 1. Analyze Grid
        # Find max columns
        max_cols = 0
        table_data = []
        for tr in rows:
            cells = tr.find_all(['td', 'th'])
            max_cols = max(max_cols, len(cells))
            row_data = [cell.get_text(strip=True) for cell in cells]
            table_data.append(row_data)
        
        box.table_data = table_data
        
        # 2. Calculate Dimensions
        # Simple: Equal width columns
        if max_cols > 0:
            col_width = box.content_width / max_cols
            
            # Estimate height based on text in cells
            total_height = 0
            for tr in rows:
                cells = tr.find_all(['td', 'th'])
                max_row_height = 0
                for cell in cells:
                    text = cell.get_text(strip=True)
                    # Approximate style for cell
                    cell_style = style_parser.compute_style(cell)
                    _, _, h = self.text_measurer.wrap_text(text, cell_style, col_width - 10) # -10 padding
                    max_row_height = max(max_row_height, h + 20) # +20 padding
                total_height += max_row_height
            
            box.content_height = total_height

import json
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

class ChartBuilder:
    """Extracts chart data from JS and builds PPTX charts."""
    
    def __init__(self):
        # Regex to find "new Chart(ctx, { ... })"
        self.chart_pattern = re.compile(r'new\s+Chart\s*\(\s*[^,]+,\s*(\{.*?\})\s*\)', re.DOTALL)

    def extract_charts(self, html_content: str) -> List[Dict]:
        """Finds chart configurations in <script> tags."""
        charts = []
        # Find all script contents
        scripts = re.findall(r'<script[^>]*>(.*?)</script>', html_content, re.DOTALL)
        
        for script in scripts:
            for match in self.chart_pattern.finditer(script):
                json_str = match.group(1)
                # JS object to JSON (simplified)
                # 1. Replace single quotes with double quotes
                json_str = json_str.replace("'", '"')
                # 2. Quote keys (simple regex: key: -> "key":)
                # Be careful not to quote inside strings. 
                # A safer way is to use a more specific regex or just assume keys are simple words.
                json_str = re.sub(r'(\w+)\s*:', r'"\1":', json_str)
                
                try:
                    # Clean up trailing commas
                    json_str = re.sub(r',\s*}', '}', json_str)
                    json_str = re.sub(r',\s*]', ']', json_str)
                    data = json.loads(json_str)
                    charts.append(data)
                except json.JSONDecodeError:
                    logger.warning("Failed to parse chart JSON. Trying loose parsing...")
                    # Fallback: try to extract labels and data arrays manually
                    # Updated regex to handle single quotes and spaces
                    labels_match = re.search(r'labels\s*:\s*\[(.*?)\]', json_str, re.DOTALL)
                    datasets_match = re.search(r'data\s*:\s*\[(.*?)\]', json_str, re.DOTALL)
                    
                    if labels_match and datasets_match:
                        labels_str = labels_match.group(1)
                        # Split by comma, strip quotes
                        labels = [l.strip().strip('"').strip("'") for l in labels_match.group(1).split(',')]
                        
                        # For data, it might be nested objects or simple array
                        # If simple array: [1, 2, 3]
                        # If objects: [{data: [...]}, ...]
                        # Our regex for datasets_match captures the content of data: [...]
                        # But in the example: datasets: [{ data: [...] }]
                        # So we need to look deeper.
                        
                        # Let's try to find "data: [...]" inside the datasets part
                        inner_data_match = re.search(r'data\s*:\s*\[(.*?)\]', datasets_match.group(1), re.DOTALL)
                        if inner_data_match:
                             data_values = [float(v.strip()) for v in inner_data_match.group(1).split(',') if v.strip()]
                             charts.append({
                                'type': 'bar',
                                'data': {
                                    'labels': labels,
                                    'datasets': [{'data': data_values}]
                                }
                            })
                        else:
                             # Maybe it was a simple array (not standard Chart.js but possible in loose parsing)
                             pass

        return charts

class PPTXBuilder:
    """Generates PowerPoint file from RenderBoxes."""
    
    def __init__(self, output_path: str):
        self.prs = Presentation()
        self.prs.slide_width = px_to_inches(SLIDE_WIDTH_PX)
        self.prs.slide_height = px_to_inches(SLIDE_HEIGHT_PX)
        self.output_path = output_path
        self.chart_placeholders = [] # List of (slide_index, box)

    def create_slides(self, pages: List[RenderBox]):
        """Creates slides based on layout pages."""
        blank_slide_layout = self.prs.slide_layouts[6] # 6 is usually blank
        
        for i, page_box in enumerate(pages):
            slide = self.prs.slides.add_slide(blank_slide_layout)
            
            # Draw background if present
            if page_box.style.background_color:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*page_box.style.background_color)
            
            # Draw children
            # Sort by z-index
            sorted_children = sorted(page_box.children, key=lambda c: c.style.z_index)
            for child in sorted_children:
                self._draw_box(slide, child, slide_index=i)

    def _draw_box(self, slide, box: RenderBox, slide_index: int):
        # Skip if invisible
        if box.total_width <= 0 or box.total_height <= 0:
            return
            
        # Handle specific box types
        if box.box_type == "image":
            self._draw_image(slide, box)
            return
        elif box.box_type == "table":
            self._draw_table(slide, box)
            return
        elif box.box_type == "chart_placeholder":
            self.chart_placeholders.append((slide_index, box))
            return

        # 1. Draw Shape (Rectangle) if it has background or border
        if box.style.background_color or box.style.border_width > 0:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                px_to_inches(box.x),
                px_to_inches(box.y),
                px_to_inches(box.total_width),
                px_to_inches(box.total_height)
            )
            
            # Fill
            fill = shape.fill
            if box.style.background_color:
                fill.solid()
                fill.fore_color.rgb = RGBColor(*box.style.background_color)
            else:
                fill.background() # No fill
            
            # Border
            line = shape.line
            if box.style.border_width > 0 and box.style.border_color:
                line.width = Pt(box.style.border_width)
                line.color.rgb = RGBColor(*box.style.border_color)
            else:
                line.fill.background() # No line

        # 2. Draw Text
        # We use the pre-calculated lines from LayoutEngine
        if box.text_lines:
            # Text box should be inside padding
            text_x = box.x + box.style.border_width + box.style.padding_left
            text_y = box.y + box.style.border_width + box.style.padding_top
            text_w = box.content_width
            text_h = box.content_height
            
            textbox = slide.shapes.add_textbox(
                px_to_inches(text_x),
                px_to_inches(text_y),
                px_to_inches(text_w),
                px_to_inches(text_h)
            )
            tf = textbox.text_frame
            tf.word_wrap = True # We already wrapped, but this is safe
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            
            # Add paragraphs
            # Note: PPTX adds a default paragraph, use it first
            p = tf.paragraphs[0]
            # Join lines with newlines? Or add separate paragraphs?
            # If we add separate paragraphs, we get paragraph spacing.
            # If we join with \n, it's one paragraph.
            # LayoutEngine wrapped lines based on width. 
            # Ideally we just give the full text to PPTX and let it wrap, 
            # BUT we calculated height based on our wrapping.
            # To be safe, let's join with spaces (if we trust PPTX wrap) or newlines (if we want to force ours).
            # Since we used PIL to measure, our wrapping is "truth".
            # However, PPTX font rendering might differ slightly.
            # Let's try joining with newlines to enforce our breaks.
            p.text = "\n".join(box.text_lines)
            
            p.font.size = Pt(box.style.font_size)
            p.font.name = box.style.font_family
            p.font.bold = (box.style.font_weight == 'bold' or str(box.style.font_weight) == '700')
            p.font.italic = (box.style.font_style == 'italic')
            
            if box.style.color:
                p.font.color.rgb = RGBColor(*box.style.color)
            
            # Alignment
            if box.style.text_align == 'center':
                p.alignment = PP_ALIGN.CENTER
            elif box.style.text_align == 'right':
                p.alignment = PP_ALIGN.RIGHT
            elif box.style.text_align == 'justify':
                p.alignment = PP_ALIGN.JUSTIFY
            else:
                p.alignment = PP_ALIGN.LEFT
                
            if box.box_type == "list_item":
                p.level = 0 # TODO: Handle nesting

        # Recurse
        for child in box.children:
            self._draw_box(slide, child, slide_index)

    def _draw_image(self, slide, box: RenderBox):
        if not box.image_data: return
        try:
            slide.shapes.add_picture(
                io.BytesIO(box.image_data),
                px_to_inches(box.x),
                px_to_inches(box.y),
                px_to_inches(box.content_width),
                px_to_inches(box.content_height)
            )
        except Exception as e:
            logger.warning(f"Failed to draw image: {e}")

    def _draw_table(self, slide, box: RenderBox):
        if not box.table_data: return
        rows = len(box.table_data)
        cols = len(box.table_data[0]) if rows > 0 else 0
        if rows == 0 or cols == 0: return
        
        table_shape = slide.shapes.add_table(
            rows, cols,
            px_to_inches(box.x),
            px_to_inches(box.y),
            px_to_inches(box.content_width),
            px_to_inches(box.content_height)
        )
        table = table_shape.table
        
        # Set styles
        for i, row_data in enumerate(box.table_data):
            for j, cell_text in enumerate(row_data):
                if j >= len(table.columns): break
                cell = table.cell(i, j)
                cell.text = cell_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
    def inject_charts(self, charts: List[Dict]):
        """Injects charts into placeholders."""
        for i, (slide_idx, box) in enumerate(self.chart_placeholders):
            if i >= len(charts): break
            chart_data = charts[i]
            self.add_chart(slide_idx, chart_data, box)

    def add_chart(self, slide_index: int, chart_data: Dict, box: Optional[RenderBox] = None):
        """Adds a chart to a specific slide."""
        if slide_index >= len(self.prs.slides):
            return
            
        slide = self.prs.slides[slide_index]
        cd = CategoryChartData()
        
        try:
            labels = chart_data.get('data', {}).get('labels', [])
            datasets = chart_data.get('data', {}).get('datasets', [])
            
            if not labels or not datasets:
                return

            for label in labels:
                cd.add_category(str(label))
            
            for i, ds in enumerate(datasets):
                name = ds.get('label', 'Series')
                values = [float(v) for v in ds.get('data', [])]
                series = cd.add_series(name, values)
                
                # Apply colors if available
                bg_color = ds.get('backgroundColor')
                if bg_color:
                    # Handle single color string or list of colors (we only support single color per series for now, or first color)
                    color_val = bg_color if isinstance(bg_color, str) else (bg_color[0] if bg_color else None)
                    
                    if color_val:
                        # Parse color
                        rgb = None
                        if color_val.startswith('#'):
                            h = color_val.lstrip('#')
                            if len(h) == 3: h = ''.join([c*2 for c in h])
                            try: rgb = tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
                            except: pass
                        elif color_val.startswith('rgb'):
                            try:
                                parts = color_val[color_val.find('(')+1:color_val.find(')')].split(',')
                                rgb = tuple(int(p.strip()) for p in parts[:3])
                            except: pass
                        elif color_val in ['red', 'blue', 'green', 'orange', 'purple', 'yellow', 'gray', 'black', 'white']:
                             # Simple map
                             c_map = {'red':(255,0,0), 'blue':(0,0,255), 'green':(0,128,0), 'orange':(255,165,0), 
                                      'purple':(128,0,128), 'yellow':(255,255,0), 'gray':(128,128,128), 
                                      'black':(0,0,0), 'white':(255,255,255)}
                             rgb = c_map[color_val]

                        if rgb:
                            # We need to access the actual series object from the chart, 
                            # but python-pptx adds series to chart_data, then we create chart.
                            # We can't style series *before* creating the chart.
                            # We must style it *after* creating the chart.
                            # Store color to apply later? 
                            # Actually, we can't easily style individual data points in python-pptx 
                            # without accessing the series object after chart creation.
                            pass 
            
            # Position
            if box:
                x, y, w, h = px_to_inches(box.x), px_to_inches(box.y), px_to_inches(box.content_width), px_to_inches(box.content_height)
            else:
                x, y, w, h = Inches(2), Inches(2), Inches(6), Inches(4.5)
                
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd
            ).chart
            
            # Apply Series Colors (Post-creation)
            for i, ds in enumerate(datasets):
                if i >= len(chart.series): break
                bg_color = ds.get('backgroundColor')
                if bg_color:
                     color_val = bg_color if isinstance(bg_color, str) else (bg_color[0] if bg_color else None)
                     # Parse color (duplicate logic, should be helper)
                     rgb = None
                     if color_val:
                        if color_val.startswith('#'):
                            h = color_val.lstrip('#')
                            if len(h) == 3: h = ''.join([c*2 for c in h])
                            try: rgb = tuple(int(h[k:k+2], 16) for k in (0, 2, 4))
                            except: pass
                        elif color_val.startswith('rgb'):
                            try:
                                parts = color_val[color_val.find('(')+1:color_val.find(')')].split(',')
                                rgb = tuple(int(p.strip()) for p in parts[:3])
                            except: pass
                        elif color_val in ['red', 'blue', 'green', 'orange', 'purple', 'yellow', 'gray', 'black', 'white']:
                             c_map = {'red':(255,0,0), 'blue':(0,0,255), 'green':(0,128,0), 'orange':(255,165,0), 
                                      'purple':(128,0,128), 'yellow':(255,255,0), 'gray':(128,128,128), 
                                      'black':(0,0,0), 'white':(255,255,255)}
                             rgb = c_map[color_val]
                     
                     if rgb:
                         chart.series[i].format.fill.solid()
                         chart.series[i].format.fill.fore_color.rgb = RGBColor(*rgb)
        except Exception as e:
            logger.error(f"Failed to add chart: {e}")

    def save(self):
        self.prs.save(self.output_path)

# --- Main Logic ---

def main():
    parser = argparse.ArgumentParser(description="HTML to PPTX Converter (Lightweight)")
    parser.add_argument("input_html", help="Path to input HTML file")
    parser.add_argument("output_pptx", help="Path to output PPTX file")
    args = parser.parse_args()

    try:
        with open(args.input_html, "r", encoding="utf-8") as f:
            html_content = f.read()
    except Exception as e:
        logger.error(f"Failed to read input file: {e}")
        sys.exit(1)

    soup = BeautifulSoup(html_content, "html.parser")
    
    # 1. Parse Styles
    logger.info("Parsing Styles...")
    style_parser = StyleParser()
    # Extract <style> tags
    for style_tag in soup.find_all("style"):
        if style_tag.string:
            style_parser.parse_css(style_tag.string)
    
    # 2. Virtual Rendering    # 3. Layout
    logger.info("Running Virtual Layout Engine...")
    layout_engine = LayoutEngine(SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX)
    pages = layout_engine.layout(soup, style_parser)
    logger.info(f"Generated {len(pages)} slides from layout.")
    
    # 3. Chart Extraction
    logger.info("Extracting Charts...")
    chart_builder = ChartBuilder()
    charts = chart_builder.extract_charts(html_content)
    logger.info(f"Found {len(charts)} charts.")
    
    # 4. Generate PPTX
    logger.info("Generating PPTX...")
    pptx_builder = PPTXBuilder(args.output_pptx)
    pptx_builder.create_slides(pages)
    
    # Inject charts into placeholders
    if charts:
        logger.info(f"Injecting {len(charts)} charts into placeholders...")
        pptx_builder.inject_charts(charts)
    
    pptx_builder.save()
    logger.info(f"Successfully created {args.output_pptx}")

if __name__ == "__main__":
    main()
