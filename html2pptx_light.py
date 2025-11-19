#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML slides -> editable PPTX converter (browser-free edition).

Dependencies (install via pip):
  pip install beautifulsoup4 python-pptx pillow

Usage:
  python html2pptx_light.py input.html output.pptx

This variant mirrors the authoring-friendly `html_ppt.py` workflow but only uses
BeautifulSoup + python-pptx (+ Pillow for colors/images). All DOM parsing,
layout recovery, and block rendering are handled locally so no Playwright runtime
is required while keeping the editable PPTX feature-set intact.
"""

from __future__ import annotations

import argparse
import math
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageColor
except ImportError:
    Image = None
    ImageColor = None

# Slide canvas size used internally (px). Actual PPTX is scaled proportionally.
SLIDE_REF_WIDTH = 1920
SLIDE_REF_HEIGHT = 1080
DEFAULT_PADDING_X = 140
DEFAULT_PADDING_Y = 120
FLOW_GAP = 40
DEFAULT_BLOCK_HEIGHT = 140
DEFAULT_LIST_INDENT = 60

# Default font sizes (px) per tag for fallback cases.
DEFAULT_FONT_SIZE = 28
TAG_FONT_SIZE = {
    "h1": 64,
    "h2": 48,
    "h3": 36,
    "h4": 30,
    "p": 28,
    "li": 26,
}

CSS_COMMENT_RE = re.compile(r"/\*.*?\*/", re.S)
SIMPLE_SELECTOR_TOKEN_RE = re.compile(r"([#.]?[\w-]+|\*)")
WHITESPACE_RE = re.compile(r"\s+")
GRID_REPEAT_RE = re.compile(r"repeat\((\d+),\s*([^)]+)\)")
CSS_COLOR_KEYWORDS = {
    "white": (255, 255, 255),
    "black": (0, 0, 0),
    "red": (255, 0, 0),
    "green": (0, 128, 0),
    "lime": (0, 255, 0),
    "blue": (0, 0, 255),
    "yellow": (255, 255, 0),
    "orange": (255, 165, 0),
    "purple": (128, 0, 128),
    "gray": (128, 128, 128),
    "grey": (128, 128, 128),
    "silver": (192, 192, 192),
    "navy": (0, 0, 128),
    "teal": (0, 128, 128),
    "cyan": (0, 255, 255),
    "aqua": (0, 255, 255),
    "magenta": (255, 0, 255),
    "fuchsia": (255, 0, 255),
    "maroon": (128, 0, 0),
    "olive": (128, 128, 0),
}



def file_or_url_to_uri(path_or_url: str) -> str:
    p = str(path_or_url)
    if p.startswith("http://") or p.startswith("https://"):
        return p
    return Path(p).resolve().as_uri()


@dataclass
class LayoutBox:
    left: Optional[float] = None
    top: Optional[float] = None
    width: Optional[float] = None
    height: Optional[float] = None


@dataclass
class TextRun:
    text: str
    font_size: Optional[float] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    color: Optional[str] = None


@dataclass
class TableCell:
    text: str
    background_color: Optional[str] = None
    border_color: Optional[str] = None
    border_width: Optional[float] = None
    text_style: Dict[str, Any] = field(default_factory=dict)
    vertical_align: Optional[str] = None


@dataclass
class Block:
    kind: str  # text, list, table, image, shape, polyline, circle, ellipse
    text: str = ""
    runs: List[TextRun] = field(default_factory=list)
    items: List[str] = field(default_factory=list)
    numbered: bool = False
    table: List[List[str]] = field(default_factory=list)
    table_cells: List[List[TableCell]] = field(default_factory=list)
    image_path: Optional[Path] = None
    image_alt: str = ""
    shape_style: Dict[str, Any] = field(default_factory=dict)
    text_style: Dict[str, Any] = field(default_factory=dict)
    vector_data: Dict[str, Any] = field(default_factory=dict)
    layout: LayoutBox = field(default_factory=LayoutBox)
    z_index: int = 0
    order: int = 0


@dataclass
class SlideModel:
    title: Optional[str]
    background_color: Optional[str]
    blocks: List[Block] = field(default_factory=list)
    layout_constraints: List["LayoutConstraint"] = field(default_factory=list)
    canvas_width: float = SLIDE_REF_WIDTH
    canvas_height: float = SLIDE_REF_HEIGHT
    scale: Optional[float] = None
    offset_x: float = 0.0
    offset_y: float = 0.0


@dataclass
class LayoutSlot:
    element_id: int
    blocks: List[Block] = field(default_factory=list)


@dataclass
class LayoutConstraint:
    kind: str  # grid or flex
    parent_style: Dict[str, str]
    parent_tag: Tag
    depth: int
    slots: List[LayoutSlot] = field(default_factory=list)
    _slot_map: Dict[int, LayoutSlot] = field(default_factory=dict, repr=False)


def compute_dom_depth(element: Tag) -> int:
    depth = 0
    current = element
    while isinstance(current, Tag):
        parent = current.parent
        if not isinstance(parent, Tag):
            break
        depth += 1
        current = parent
    return depth


def find_direct_child_for_parent(element: Tag, ancestor: Tag) -> Optional[Tag]:
    child = element
    parent = element.parent
    while parent and parent is not ancestor:
        if not isinstance(parent, Tag):
            return None
        child = parent
        parent = parent.parent
    return child if parent is ancestor else None


def register_layout_constraints(
    element: Tag,
    block: Block,
    resolver: "StyleResolver",
    constraints: Dict[int, LayoutConstraint],
) -> None:
    parent = element.parent
    while isinstance(parent, Tag):
        style = resolver.get_style(parent)
        display = style.get("display", "").lower()
        if display in {"grid", "flex"}:
            key = id(parent)
            constraint = constraints.get(key)
            if not constraint:
                constraint = LayoutConstraint(
                    kind=display,
                    parent_style=dict(style),
                    parent_tag=parent,
                    depth=compute_dom_depth(parent),
                )
                constraints[key] = constraint
            direct_child = find_direct_child_for_parent(element, parent)
            if direct_child is not None:
                slot_id = id(direct_child)
                slot = constraint._slot_map.get(slot_id)
                if not slot:
                    slot = LayoutSlot(element_id=slot_id)
                    constraint._slot_map[slot_id] = slot
                    constraint.slots.append(slot)
                slot.blocks.append(block)
        parent = parent.parent if isinstance(parent.parent, Tag) else None


def parse_gap_values(style: Dict[str, str]) -> Tuple[float, float]:
    row_gap = parse_length(style.get("row-gap"))
    col_gap = parse_length(style.get("column-gap"))
    gap_val = style.get("gap")
    if gap_val:
        parts = gap_val.split()
        if len(parts) == 1:
            gap = parse_length(parts[0])
            if gap is not None:
                if row_gap is None:
                    row_gap = gap
                if col_gap is None:
                    col_gap = gap
        elif len(parts) >= 2:
            first = parse_length(parts[0])
            second = parse_length(parts[1])
            if row_gap is None:
                row_gap = first
            if col_gap is None:
                col_gap = second
    if row_gap is None:
        row_gap = FLOW_GAP
    if col_gap is None:
        col_gap = FLOW_GAP
    return row_gap, col_gap


def expand_grid_template_tokens(template: str) -> List[str]:
    tokens: List[str] = []
    if not template:
        return tokens
    buffer = ""
    depth = 0
    for ch in template:
        if ch == "(":
            depth += 1
        elif ch == ")":
            depth = max(0, depth - 1)
        if ch.isspace() and depth == 0:
            if buffer:
                tokens.append(buffer.strip())
                buffer = ""
        else:
            buffer += ch
    if buffer:
        tokens.append(buffer.strip())
    expanded: List[str] = []
    for token in tokens:
        if not token:
            continue
        match = GRID_REPEAT_RE.match(token)
        if match:
            count = int(match.group(1))
            inner = match.group(2).strip()
            for _ in range(count):
                expanded.append(inner)
        else:
            expanded.append(token)
    return expanded


def compute_grid_column_widths(
    container_width: float,
    template: Optional[str],
    column_gap: float,
    block_count: int,
) -> List[float]:
    tokens = expand_grid_template_tokens(template or "")
    if not tokens:
        count = max(1, min(block_count, 4))
        available = max(container_width - column_gap * (count - 1), 100.0)
        width = available / count
        return [width for _ in range(count)]
    column_defs: List[Tuple[str, float]] = []
    px_total = 0.0
    fr_total = 0.0
    for token in tokens:
        norm = token.lower()
        if norm.endswith("fr"):
            try:
                value = float(norm[:-2] or "1")
            except ValueError:
                value = 1.0
            fr_total += value
            column_defs.append(("fr", value))
        else:
            length = parse_length(token)
            if length is None:
                fr_total += 1.0
                column_defs.append(("fr", 1.0))
            else:
                px_total += length
                column_defs.append(("px", length))
    count = len(column_defs)
    available = max(container_width - column_gap * (count - 1), 50.0)
    widths: List[float] = []
    if px_total > available and px_total > 0:
        scale = available / px_total
        column_defs = [
            (kind, value * scale) if kind == "px" else (kind, value) for kind, value in column_defs
        ]
        px_total = available
    remaining = max(available - px_total, 0.0)
    for kind, value in column_defs:
        if kind == "px":
            widths.append(value)
        else:
            portion = remaining / fr_total if fr_total else (remaining / count if count else remaining)
            widths.append(portion * value if fr_total else portion)
    return widths


def compute_slot_metrics(slot: LayoutSlot) -> Tuple[float, float, float, float]:
    min_left = None
    min_top = None
    max_right = None
    max_bottom = None
    for block in slot.blocks:
        left = block.layout.left if block.layout.left is not None else DEFAULT_PADDING_X
        top = block.layout.top if block.layout.top is not None else DEFAULT_PADDING_Y
        width = block.layout.width if block.layout.width is not None else SLIDE_REF_WIDTH - 2 * DEFAULT_PADDING_X
        height = block.layout.height if block.layout.height is not None else estimate_block_height(block)
        min_left = left if min_left is None else min(min_left, left)
        min_top = top if min_top is None else min(min_top, top)
        right = left + width
        bottom = top + height
        max_right = right if max_right is None else max(max_right, right)
        max_bottom = bottom if max_bottom is None else max(max_bottom, bottom)
    min_left = min_left if min_left is not None else DEFAULT_PADDING_X
    min_top = min_top if min_top is not None else DEFAULT_PADDING_Y
    width_span = (max_right - min_left) if max_right is not None else SLIDE_REF_WIDTH - 2 * DEFAULT_PADDING_X
    height_span = (max_bottom - min_top) if max_bottom is not None else DEFAULT_BLOCK_HEIGHT
    return float(min_left), float(min_top), float(max(width_span, 10.0)), float(max(height_span, 10.0))


def compute_container_box(constraint: LayoutConstraint) -> Tuple[float, float, float]:
    layout = extract_layout(constraint.parent_style, constraint.parent_tag)
    block_lefts = [
        block.layout.left
        for slot in constraint.slots
        for block in slot.blocks
        if block.layout.left is not None
    ]
    block_rights = [
        (block.layout.left + block.layout.width)
        for slot in constraint.slots
        for block in slot.blocks
        if block.layout.left is not None and block.layout.width is not None
    ]
    block_tops = [
        block.layout.top
        for slot in constraint.slots
        for block in slot.blocks
        if block.layout.top is not None
    ]
    left = layout.left if layout.left is not None else (min(block_lefts) if block_lefts else DEFAULT_PADDING_X)
    top = layout.top if layout.top is not None else (min(block_tops) if block_tops else DEFAULT_PADDING_Y)
    if layout.width is not None:
        width = layout.width
    elif block_rights:
        width = max(block_rights) - left
    else:
        width = SLIDE_REF_WIDTH - 2 * DEFAULT_PADDING_X
    return float(left), float(top), float(max(width, 100.0))


def px_to_pt(px: float) -> float:
    return float(px) * 0.75


def css_color_to_rgb_tuple(color_str: Optional[str]) -> Optional[Tuple[int, int, int]]:
    if not color_str:
        return None
    s = color_str.strip()
    if not s:
        return None
    try:
        if s.startswith("#"):
            if len(s) == 4:  # #RGB
                r = int(s[1], 16) * 17
                g = int(s[2], 16) * 17
                b = int(s[3], 16) * 17
                return (r, g, b)
            if len(s) == 7:
                return (int(s[1:3], 16), int(s[3:5], 16), int(s[5:7], 16))
        if s.lower().startswith("rgb"):
            nums = s[s.find("(") + 1 : s.find(")")].split(",")
            r, g, b = [int(float(v.strip())) for v in nums[:3]]
            return (r, g, b)
    except Exception:
        return None
    keyword = CSS_COLOR_KEYWORDS.get(s.lower())
    if keyword:
        return keyword
    if ImageColor:
        try:
            return tuple(ImageColor.getrgb(s))
        except Exception:
            return None
    return None


def css_is_transparent(color_str: Optional[str]) -> bool:
    if not color_str:
        return True
    s = color_str.strip().lower()
    if not s:
        return True
    if s == "transparent":
        return True
    if s.startswith("rgba"):
        try:
            nums = s[s.find("(") + 1 : s.find(")")].split(",")
            if len(nums) >= 4:
                alpha = float(nums[3])
                return alpha <= 0.0
        except Exception:
            return False
    return False


def parse_length(value: Optional[str], reference: Optional[float] = None) -> Optional[float]:
    if value is None:
        return None
    s = value.strip()
    if not s or s in {"auto", "initial", "inherit"}:
        return None
    if s.endswith("px"):
        return float(s[:-2])
    if s.endswith("%") and reference is not None:
        try:
            return float(s[:-1]) / 100.0 * reference
        except ValueError:
            return None
    if s.endswith("em") or s.endswith("rem"):
        try:
            return float(s[:-2]) * DEFAULT_FONT_SIZE
        except ValueError:
            return None
    try:
        return float(s)
    except ValueError:
        return None


def parse_font_size(value: Optional[str], tag: Optional[str] = None) -> float:
    size = parse_length(value)
    if size:
        return size
    if tag and tag.lower() in TAG_FONT_SIZE:
        return TAG_FONT_SIZE[tag.lower()]
    return DEFAULT_FONT_SIZE


def normalize_whitespace(text: str) -> str:
    return WHITESPACE_RE.sub(" ", text).strip()


def parse_declarations(text: str) -> Dict[str, str]:
    result: Dict[str, str] = {}
    for part in text.split(";"):
        if ":" not in part:
            continue
        name, val = part.split(":", 1)
        name = name.strip().lower()
        if not name:
            continue
        result[name] = val.strip()
    return result


@dataclass
class StyleRule:
    tag: Optional[str]
    element_id: Optional[str]
    classes: Tuple[str, ...]
    declarations: Dict[str, str]

    def matches(self, tag: str, classes: Iterable[str], element_id: Optional[str]) -> bool:
        class_set = set(classes or [])
        if self.tag and self.tag != "*" and self.tag != tag:
            return False
        if self.element_id and self.element_id != element_id:
            return False
        if self.classes and not set(self.classes).issubset(class_set):
            return False
        return True


def parse_simple_selector(selector: str) -> Optional[StyleRule]:
    selector = selector.strip()
    if not selector or " " in selector or ">" in selector:
        return None
    tokens = SIMPLE_SELECTOR_TOKEN_RE.findall(selector)
    if not tokens:
        return None
    tag: Optional[str] = None
    el_id: Optional[str] = None
    classes: List[str] = []
    for token in tokens:
        if token == "*":
            tag = "*"
        elif token.startswith("#"):
            el_id = token[1:]
        elif token.startswith("."):
            classes.append(token[1:])
        else:
            tag = token.lower()
    return StyleRule(tag=tag, element_id=el_id, classes=tuple(classes), declarations={})


class StyleResolver:
    """Very small CSS resolver (supports tag/id/class selectors without combinators)."""

    def __init__(self, soup: BeautifulSoup):
        self.rules: List[StyleRule] = []
        for style_tag in soup.select("style"):
            self._consume_stylesheet(style_tag.string or "")

    def _consume_stylesheet(self, css_text: str) -> None:
        cleaned = CSS_COMMENT_RE.sub("", css_text)
        for raw_rule in cleaned.split("}"):
            if "{" not in raw_rule:
                continue
            selector_text, body = raw_rule.split("{", 1)
            declarations = parse_declarations(body)
            if not declarations:
                continue
            for selector in selector_text.split(","):
                rule = parse_simple_selector(selector)
                if not rule:
                    continue
                rule.declarations = declarations.copy()
                self.rules.append(rule)

    def get_style(self, element: Tag) -> Dict[str, str]:
        tag = (element.name or "").lower()
        classes = element.get("class", [])
        element_id = element.get("id")
        style: Dict[str, str] = {}
        for rule in self.rules:
            if rule.matches(tag, classes, element_id):
                style.update(rule.declarations)
        inline = element.get("style")
        if inline:
            style.update(parse_declarations(inline))
        return style


def build_text_style(tag: str, style: Dict[str, str]) -> Dict[str, Any]:
    result: Dict[str, Any] = {}
    result["font_size"] = parse_font_size(style.get("font-size"), tag)
    weight = style.get("font-weight", "").lower()
    if weight in {"bold", "bolder"}:
        result["bold"] = True
    elif weight.isdigit():
        result["bold"] = int(weight) >= 600
    italic = style.get("font-style", "").lower()
    if italic:
        result["italic"] = italic == "italic"
    color = style.get("color")
    if color:
        result["color"] = color
    align = style.get("text-align", "").lower()
    if align in {"left", "center", "right", "justify"}:
        result["align"] = align
    return result


def css_weight_is_bold(weight: Optional[str]) -> bool:
    if weight is None:
        return False
    s = str(weight).strip().lower()
    if not s:
        return False
    if s.isdigit():
        try:
            return int(s) >= 600
        except ValueError:
            return False
    return s in {"bold", "bolder", "600", "700", "800", "900"}


def apply_text_style(base: Dict[str, Any], style: Dict[str, str], tag: Optional[str] = None) -> Dict[str, Any]:
    new_style = dict(base)
    size = parse_length(style.get("font-size"))
    if size:
        new_style["font_size"] = size
    weight = style.get("font-weight", "").lower()
    if weight in {"bold", "bolder"}:
        new_style["bold"] = True
    elif weight.isdigit():
        new_style["bold"] = int(weight) >= 600
    italic = style.get("font-style", "").lower()
    if italic:
        new_style["italic"] = italic == "italic"
    color = style.get("color")
    if color:
        new_style["color"] = color
    align = style.get("text-align", "").lower()
    if align in {"left", "center", "right", "justify"}:
        new_style["align"] = align
    if tag in {"strong", "b"}:
        new_style["bold"] = True
    if tag in {"em", "i"}:
        new_style["italic"] = True
    return new_style


def extract_layout(style: Dict[str, str], element: Optional[Tag] = None) -> LayoutBox:
    box = LayoutBox()
    attr_width = element.get("width") if element and element.has_attr("width") else None
    attr_height = element.get("height") if element and element.has_attr("height") else None
    box.width = parse_length(attr_width, SLIDE_REF_WIDTH)
    box.height = parse_length(attr_height, SLIDE_REF_HEIGHT)
    width_style = style.get("width")
    height_style = style.get("height")
    if width_style:
        box.width = parse_length(width_style, SLIDE_REF_WIDTH)
    if height_style:
        box.height = parse_length(height_style, SLIDE_REF_HEIGHT)
    box.left = parse_length(style.get("left"), SLIDE_REF_WIDTH)
    box.top = parse_length(style.get("top"), SLIDE_REF_HEIGHT)
    return box


def merge_runs(runs: List[TextRun]) -> List[TextRun]:
    merged: List[TextRun] = []
    for run in runs:
        if not run.text:
            continue
        if merged:
            prev = merged[-1]
            if (
                prev.font_size == run.font_size
                and prev.bold == run.bold
                and prev.italic == run.italic
                and prev.color == run.color
            ):
                prev.text += run.text
                continue
        merged.append(run)
    return merged


def extract_text_runs(element: Tag, resolver: StyleResolver, base_style: Dict[str, Any]) -> List[TextRun]:
    runs: List[TextRun] = []

    def walk(node: Any, current_style: Dict[str, Any]) -> None:
        if isinstance(node, NavigableString):
            text = str(node)
            normalized = WHITESPACE_RE.sub(" ", text)
            if normalized.strip():
                runs.append(
                    TextRun(
                        text=normalized,
                        font_size=current_style.get("font_size"),
                        bold=current_style.get("bold"),
                        italic=current_style.get("italic"),
                        color=current_style.get("color"),
                    )
                )
            return
        if not isinstance(node, Tag):
            return
        if node.name == "br":
            runs.append(
                TextRun(
                    text="\n",
                    font_size=current_style.get("font_size"),
                    bold=current_style.get("bold"),
                    italic=current_style.get("italic"),
                    color=current_style.get("color"),
                )
            )
            return
        node_style = resolver.get_style(node)
        next_style = apply_text_style(current_style, node_style, tag=node.name)
        for child in node.children:
            walk(child, next_style)

    for child in element.children:
        walk(child, dict(base_style))
    if not runs:
        text = normalize_whitespace(element.get_text(" ", strip=True))
        if text:
            runs.append(
                TextRun(
                    text=text,
                    font_size=base_style.get("font_size"),
                    bold=base_style.get("bold"),
                    italic=base_style.get("italic"),
                    color=base_style.get("color"),
                )
            )
    return merge_runs(runs)


def detect_shape_style(style: Dict[str, str]) -> Dict[str, Any]:
    fill_color = style.get("background-color") or style.get("background")
    border_value = style.get("border", "")
    border_width = parse_length(style.get("border-width"))
    border_color = style.get("border-color")
    if border_value:
        for token in border_value.split():
            length = parse_length(token)
            if length is not None:
                border_width = length
            elif token.startswith("#") or token.startswith("rgb"):
                border_color = token
    border_radius = parse_length(style.get("border-radius"))
    if not any([fill_color, border_color, border_width, border_radius]):
        return {}
    return {
        "fill_color": fill_color,
        "border_color": border_color,
        "border_width": border_width,
        "border_radius": border_radius,
    }


def parse_z_index(style: Dict[str, str]) -> int:
    z_str = (style.get("z-index") or "").strip().lower()
    if not z_str or z_str == "auto":
        return 0
    try:
        return int(float(z_str))
    except ValueError:
        return 0


def has_block_children(element: Tag) -> bool:
    block_tags = {"p", "div", "section", "article", "ul", "ol", "table", "h1", "h2", "h3", "h4", "figure"}
    for child in element.find_all(block_tags, recursive=False):
        if child is not None:
            return True
    return False


def resolve_image_path(src: str, base_dir: Path) -> Optional[Path]:
    if not src:
        return None
    parsed = re.match(r"^(?P<scheme>[a-zA-Z][a-zA-Z0-9+.-]*):", src)
    if parsed:
        scheme = parsed.group("scheme").lower()
        if scheme in {"http", "https"}:
            return None  # Remote resources are not fetched.
        if scheme == "file":
            return Path(src.split("://", 1)[-1])
        if scheme == "data":
            return None
    candidate = (base_dir / src).resolve()
    return candidate if candidate.exists() else None


def element_text(element: Tag) -> str:
    return element.get_text(" ", strip=True)


def extract_blocks(
    slide_el: Tag,
    resolver: StyleResolver,
    base_dir: Path,
    constraints_map: Dict[int, LayoutConstraint],
) -> List[Block]:
    blocks: List[Block] = []
    order_counter = 0
    for element in slide_el.descendants:
        if not isinstance(element, Tag):
            continue
        if element.name in {"script", "style", "noscript"}:
            continue
        tag = element.name.lower()
        style = resolver.get_style(element)
        layout = extract_layout(style, element)
        text_style = build_text_style(tag, style)
        block: Optional[Block] = None

        if tag in {"h1", "h2", "h3", "h4", "p", "blockquote"}:
            runs = extract_text_runs(element, resolver, text_style)
            text = element_text(element)
            block = Block(kind="text", text=text, runs=runs, text_style=text_style, layout=layout)
            shape_style = detect_shape_style(style)
            if shape_style:
                block.shape_style = shape_style

        elif tag in {"ul", "ol"}:
            items = [li.get_text(" ", strip=True) for li in element.find_all("li", recursive=False)]
            items = [item for item in items if item]
            if items:
                block = Block(
                    kind="list",
                    items=items,
                    numbered=(tag == "ol"),
                    text_style=text_style,
                    layout=layout,
                )

        elif tag == "table":
            rows: List[List[str]] = []
            cell_styles: List[List[TableCell]] = []
            for tr in element.find_all("tr", recursive=False):
                row = []
                row_cells: List[TableCell] = []
                for cell in tr.find_all(["th", "td"], recursive=False):
                    cell_text = cell.get_text(" ", strip=True)
                    row.append(cell_text)
                    cell_style = resolver.get_style(cell)
                    cell_text_style = build_text_style(cell.name, cell_style)
                    row_cells.append(
                        TableCell(
                            text=cell_text,
                            background_color=cell_style.get("background-color") or cell_style.get("background"),
                            border_color=cell_style.get("border-color"),
                            border_width=parse_length(cell_style.get("border-width")),
                            text_style=cell_text_style,
                            vertical_align=cell_style.get("vertical-align"),
                        )
                    )
                if row:
                    rows.append(row)
                    cell_styles.append(row_cells)
            if rows:
                block = Block(kind="table", table=rows, table_cells=cell_styles, text_style=text_style, layout=layout)

        elif tag == "img":
            src = element.get("src")
            img_path = resolve_image_path(src, base_dir) if src else None
            block = Block(
                kind="image",
                image_path=img_path,
                image_alt=element.get("alt", ""),
                layout=layout,
            )

        elif tag in {"div", "span", "section"}:
            text_content = element_text(element)
            shape_style = detect_shape_style(style)
            if text_content and not has_block_children(element):
                runs = extract_text_runs(element, resolver, text_style)
                block = Block(kind="text", text=text_content, runs=runs, text_style=text_style, layout=layout)
                if shape_style:
                    block.shape_style = shape_style
            elif shape_style:
                block = Block(kind="shape", shape_style=shape_style, layout=layout)

        elif tag == "hr":
            block = Block(
                kind="shape",
                shape_style={"fill_color": style.get("background-color") or "#999999"},
                layout=layout,
            )

        if block:
            block.z_index = parse_z_index(style)
            block.order = order_counter
            order_counter += 1
            blocks.append(block)
            register_layout_constraints(element, block, resolver, constraints_map)
    return blocks


def estimate_block_height(block: Block) -> float:
    if block.layout.height:
        return block.layout.height
    style = block.text_style or {}
    font_size = style.get("font_size", DEFAULT_FONT_SIZE)
    if block.kind == "text":
        lines = max(1, block.text.count("\n") + 1)
        return max(80, lines * font_size * 1.4 + 20)
    if block.kind == "list":
        lines = max(1, len(block.items))
        return max(80, lines * font_size * 1.3 + 20)
    if block.kind == "table":
        rows = max(1, len(block.table))
        return rows * (font_size * 1.8)
    if block.kind == "image":
        return 240
    if block.kind == "shape":
        return 120
    return DEFAULT_BLOCK_HEIGHT


def assign_fallback_layouts(slides: List[SlideModel]) -> None:
    for slide in slides:
        slide_width = slide.canvas_width or SLIDE_REF_WIDTH
        flow_y = DEFAULT_PADDING_Y
        for block in slide.blocks:
            if block.layout.left is None:
                block.layout.left = DEFAULT_PADDING_X
            if block.layout.width is None:
                block.layout.width = max(slide_width - 2 * DEFAULT_PADDING_X, 200.0)
            block_height = block.layout.height or estimate_block_height(block)
            if block.layout.top is None:
                block.layout.top = flow_y
            block.layout.height = block_height
            flow_y = max(flow_y, block.layout.top + block.layout.height + FLOW_GAP)


def apply_layout_constraints(slides: List[SlideModel]) -> None:
    for slide in slides:
        if not slide.layout_constraints:
            continue
        constraints = sorted(slide.layout_constraints, key=lambda c: c.depth)
        for constraint in constraints:
            slots = [slot for slot in constraint.slots if slot.blocks]
            if len(slots) <= 1:
                continue
            if constraint.kind == "grid":
                apply_grid_constraint(constraint, slots)
            elif constraint.kind == "flex":
                apply_flex_constraint(constraint, slots)


def apply_grid_constraint(constraint: LayoutConstraint, slots: List[LayoutSlot]) -> None:
    container_left, container_top, container_width = compute_container_box(constraint)
    row_gap, col_gap = parse_gap_values(constraint.parent_style)
    column_widths = compute_grid_column_widths(
        container_width,
        constraint.parent_style.get("grid-template-columns"),
        col_gap,
        len(slots),
    )
    if not column_widths:
        column_widths = [container_width]
    col_count = max(1, len(column_widths))
    column_offsets: List[float] = []
    accum = 0.0
    for width in column_widths:
        column_offsets.append(accum)
        accum += width + col_gap
    row_top = container_top
    row_max_height = 0.0
    for idx, slot in enumerate(slots):
        col_index = idx % col_count
        if idx > 0 and col_index == 0:
            row_top += row_max_height + row_gap
            row_max_height = 0.0
        width = column_widths[col_index]
        slot_left = container_left + column_offsets[col_index]
        min_left, min_top, slot_width, slot_height = compute_slot_metrics(slot)
        slot_height = slot_height or DEFAULT_BLOCK_HEIGHT
        left_offset = slot_left - min_left
        top_offset = row_top - min_top
        for block in slot.blocks:
            orig_left = block.layout.left if block.layout.left is not None else min_left
            orig_top = block.layout.top if block.layout.top is not None else min_top
            orig_width = block.layout.width if block.layout.width is not None else slot_width
            block.layout.left = orig_left + left_offset
            block.layout.top = orig_top + top_offset
            block.layout.width = min(orig_width, width)
        row_max_height = max(row_max_height, slot_height)


def apply_flex_constraint(constraint: LayoutConstraint, slots: List[LayoutSlot]) -> None:
    container_left, container_top, container_width = compute_container_box(constraint)
    row_gap, col_gap = parse_gap_values(constraint.parent_style)
    direction = constraint.parent_style.get("flex-direction", "row").lower()
    metrics = [compute_slot_metrics(slot) for slot in slots]
    if direction.startswith("column"):
        top = container_top
        for slot, metric in zip(slots, metrics):
            min_left, min_top, slot_width, slot_height = metric
            slot_height = slot_height or DEFAULT_BLOCK_HEIGHT
            left_offset = container_left - min_left
            top_offset = top - min_top
            for block in slot.blocks:
                orig_left = block.layout.left if block.layout.left is not None else min_left
                orig_top = block.layout.top if block.layout.top is not None else min_top
                orig_width = block.layout.width if block.layout.width is not None else slot_width
                block.layout.left = orig_left + left_offset
                block.layout.top = orig_top + top_offset
                block.layout.width = min(orig_width, container_width)
            top += slot_height + row_gap
    else:
        width_available = max(container_width - col_gap * (len(slots) - 1), 50.0)
        width_per = width_available / len(slots) if slots else width_available
        left = container_left
        for slot, metric in zip(slots, metrics):
            min_left, min_top, slot_width, slot_height = metric
            left_offset = left - min_left
            top_offset = container_top - min_top
            for block in slot.blocks:
                orig_left = block.layout.left if block.layout.left is not None else min_left
                orig_top = block.layout.top if block.layout.top is not None else min_top
                orig_width = block.layout.width if block.layout.width is not None else slot_width
                block.layout.left = orig_left + left_offset
                block.layout.top = orig_top + top_offset
                block.layout.width = min(orig_width, width_per)
            left += width_per + col_gap


def fit_slide_content(slide: SlideModel) -> None:
    if not slide.blocks:
        return
    width = slide.canvas_width or SLIDE_REF_WIDTH
    height = slide.canvas_height or SLIDE_REF_HEIGHT
    min_left = min((block.layout.left for block in slide.blocks if block.layout.left is not None), default=0.0)
    min_top = min((block.layout.top for block in slide.blocks if block.layout.top is not None), default=0.0)
    default_width = max(width - 2 * DEFAULT_PADDING_X, 200.0)
    max_right = max(
        (
            (block.layout.left or 0.0) + (block.layout.width or default_width)
            for block in slide.blocks
        ),
        default=width,
    )
    max_bottom = max(
        (
            (block.layout.top or 0.0) + (block.layout.height or estimate_block_height(block))
            for block in slide.blocks
        ),
        default=height,
    )
    content_width = max_right - min_left
    content_height = max_bottom - min_top
    avail_width = max(width - 2 * DEFAULT_PADDING_X, 200.0)
    avail_height = max(height - 2 * DEFAULT_PADDING_Y, 200.0)
    scale = 1.0
    needs_fit = False
    if content_width > avail_width + 1:
        scale = min(scale, avail_width / content_width)
        needs_fit = True
    if content_height > avail_height + 1:
        scale = min(scale, avail_height / content_height)
        needs_fit = True
    if min_left < DEFAULT_PADDING_X - 5 or min_top < DEFAULT_PADDING_Y - 5:
        needs_fit = True
    if max_right > width - DEFAULT_PADDING_X + 5 or max_bottom > height - DEFAULT_PADDING_Y + 5:
        needs_fit = True
    if not needs_fit:
        return
    if scale <= 0:
        scale = 1.0
    for block in slide.blocks:
        left = block.layout.left or 0.0
        top = block.layout.top or 0.0
        width_px = block.layout.width or (content_width if content_width > 0 else avail_width)
        height_px = block.layout.height or (content_height if content_height > 0 else avail_height)
        block.layout.left = DEFAULT_PADDING_X + (left - min_left) * scale
        block.layout.top = DEFAULT_PADDING_Y + (top - min_top) * scale
        block.layout.width = width_px * scale
        block.layout.height = height_px * scale


def parse_html_static(input_html: str, selector: Optional[str] = None) -> List[SlideModel]:
    path = Path(input_html)
    base_dir = path.parent
    with path.open("r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
    resolver = StyleResolver(soup)
    if selector:
        slide_candidates = soup.select(selector)
    else:
        slide_candidates = soup.select(".slide")
    if not slide_candidates:
        slide_candidates = soup.select("[data-slide]")
    if not slide_candidates:
        if soup.body:
            slide_candidates = [soup.body]
        else:
            raise ValueError("HTML に <body> が見つかりません。スライドを判定できません。")
    slides: List[SlideModel] = []
    for idx, slide_el in enumerate(slide_candidates):
        style = resolver.get_style(slide_el)
        background = style.get("background-color") or style.get("background")
        title_el = slide_el.find(["h1", "h2", "h3", "h4"])
        title = title_el.get_text(" ", strip=True) if title_el else None
        constraints_map: Dict[int, LayoutConstraint] = {}
        blocks = extract_blocks(slide_el, resolver, base_dir, constraints_map)
        constraints = sorted(constraints_map.values(), key=lambda c: c.depth)
        canvas_width = (
            parse_length(slide_el.get("width")) or parse_length(style.get("width")) or SLIDE_REF_WIDTH
        )
        canvas_height = (
            parse_length(slide_el.get("height"))
            or parse_length(style.get("height"))
            or parse_length(style.get("min-height"))
            or SLIDE_REF_HEIGHT
        )
        slides.append(
            SlideModel(
                title=title,
                background_color=background,
                blocks=blocks,
                layout_constraints=constraints,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
            )
        )
    assign_fallback_layouts(slides)
    apply_layout_constraints(slides)
    for slide in slides:
        fit_slide_content(slide)
    return slides


def ensure_slide_transform(slide_model: SlideModel, prs: Presentation) -> None:
    if slide_model.scale:
        return
    canvas_width = slide_model.canvas_width or SLIDE_REF_WIDTH
    canvas_height = slide_model.canvas_height or SLIDE_REF_HEIGHT
    if canvas_width <= 0:
        canvas_width = SLIDE_REF_WIDTH
    if canvas_height <= 0:
        canvas_height = SLIDE_REF_HEIGHT
    scale_x = prs.slide_width / canvas_width
    scale_y = prs.slide_height / canvas_height
    scale = min(scale_x, scale_y)
    if scale <= 0:
        scale = scale_x or scale_y or 1.0
    slide_model.scale = scale
    slide_model.offset_x = (prs.slide_width - canvas_width * scale) / 2
    slide_model.offset_y = (prs.slide_height - canvas_height * scale) / 2


def length_to_emu(value_px: float, axis: str, prs: Presentation, slide_model: SlideModel) -> int:
    ensure_slide_transform(slide_model, prs)
    scale = slide_model.scale or 1.0
    return int(max(value_px, 0.0) * scale)


def position_to_emu(value_px: float, axis: str, prs: Presentation, slide_model: SlideModel) -> int:
    ensure_slide_transform(slide_model, prs)
    scale = slide_model.scale or 1.0
    base = slide_model.offset_x if axis == "x" else slide_model.offset_y
    return int(base + value_px * scale)


def apply_paragraph_alignment(paragraph, align: Optional[str]) -> None:
    if not align:
        return
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    paragraph.alignment = align_map.get(align, PP_ALIGN.LEFT)


def add_text_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    left = position_to_emu(block.layout.left or DEFAULT_PADDING_X, "x", prs, slide_model)
    top = position_to_emu(block.layout.top or DEFAULT_PADDING_Y, "y", prs, slide_model)
    fallback_width = max((slide_model.canvas_width or SLIDE_REF_WIDTH) - 2 * DEFAULT_PADDING_X, 200.0)
    width = length_to_emu(block.layout.width or fallback_width, "x", prs, slide_model)
    height = length_to_emu(block.layout.height or estimate_block_height(block), "y", prs, slide_model)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    base_style = block.text_style or {}
    runs = block.runs or [
        TextRun(
            text=block.text,
            font_size=base_style.get("font_size"),
            bold=base_style.get("bold"),
            italic=base_style.get("italic"),
            color=base_style.get("color"),
        )
    ]
    paragraph = tf.paragraphs[0]
    apply_paragraph_alignment(paragraph, base_style.get("align"))
    for run in runs:
        pieces = run.text.split("\n")
        for idx, piece in enumerate(pieces):
            if idx > 0:
                paragraph = tf.add_paragraph()
                apply_paragraph_alignment(paragraph, base_style.get("align"))
            ppt_run = paragraph.add_run()
            ppt_run.text = piece
            font = ppt_run.font
            size = run.font_size or base_style.get("font_size") or DEFAULT_FONT_SIZE
            font.size = Pt(px_to_pt(size))
            font.bold = run.bold if run.bold is not None else base_style.get("bold")
            font.italic = run.italic if run.italic is not None else base_style.get("italic")
            color = run.color or base_style.get("color")
            rgb = css_color_to_rgb_tuple(color)
            if rgb:
                font.color.rgb = RGBColor(*rgb)
    if block.shape_style:
        shape_style = block.shape_style
        fill_color = css_color_to_rgb_tuple(shape_style.get("fill_color"))
        if fill_color:
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*fill_color)
        else:
            box.fill.background()
        border_color = css_color_to_rgb_tuple(shape_style.get("border_color"))
        if border_color:
            line = box.line
            line.color.rgb = RGBColor(*border_color)
            width = shape_style.get("border_width") or 2
            line.width = Pt(px_to_pt(width))


def add_list_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    left = position_to_emu(
        block.layout.left or DEFAULT_PADDING_X + DEFAULT_LIST_INDENT, "x", prs, slide_model
    )
    top = position_to_emu(block.layout.top or DEFAULT_PADDING_Y, "y", prs, slide_model)
    fallback_width = max(
        (slide_model.canvas_width or SLIDE_REF_WIDTH) - 2 * DEFAULT_PADDING_X - DEFAULT_LIST_INDENT,
        200.0,
    )
    width = length_to_emu(block.layout.width or fallback_width, "x", prs, slide_model)
    height = length_to_emu(block.layout.height or estimate_block_height(block), "y", prs, slide_model)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    base_style = block.text_style or {}
    for idx, item in enumerate(block.items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        text = item.strip()
        if block.numbered:
            text = f"{idx + 1}. {text}"
        paragraph.text = text
        apply_paragraph_alignment(paragraph, base_style.get("align"))
        font = paragraph.font
        size = base_style.get("font_size", DEFAULT_FONT_SIZE)
        font.size = Pt(px_to_pt(size))
        if base_style.get("bold") is not None:
            font.bold = bool(base_style["bold"])
        if base_style.get("italic") is not None:
            font.italic = bool(base_style["italic"])
        rgb = css_color_to_rgb_tuple(base_style.get("color"))
        if rgb:
            font.color.rgb = RGBColor(*rgb)


def apply_cell_border(cell, color: Optional[str], width_px: float) -> None:
    if not color or css_is_transparent(color):
        return
    rgb = css_color_to_rgb_tuple(color)
    if not rgb:
        return
    width_pt = px_to_pt(width_px or 1.0)
    width_emu = Pt(width_pt).emu
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tcPr.find(qn(tag))
        if ln is None:
            ln = OxmlElement(tag)
            tcPr.append(ln)
        ln.set("w", str(int(width_emu)))
        solid = ln.find(qn("a:solidFill"))
        if solid is None:
            solid = OxmlElement("a:solidFill")
            ln.append(solid)
        srgb = solid.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = OxmlElement("a:srgbClr")
            solid.append(srgb)
        srgb.set("val", "{:02X}{:02X}{:02X}".format(*rgb))
        prst = ln.find(qn("a:prstDash"))
        if prst is None:
            prst = OxmlElement("a:prstDash")
            prst.set("val", "solid")
            ln.append(prst)


def add_table_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    left = position_to_emu(block.layout.left or DEFAULT_PADDING_X, "x", prs, slide_model)
    top = position_to_emu(block.layout.top or DEFAULT_PADDING_Y, "y", prs, slide_model)
    fallback_width = max((slide_model.canvas_width or SLIDE_REF_WIDTH) - 2 * DEFAULT_PADDING_X, 200.0)
    width = length_to_emu(block.layout.width or fallback_width, "x", prs, slide_model)
    height = length_to_emu(block.layout.height or estimate_block_height(block), "y", prs, slide_model)
    rows = len(block.table)
    cols = max(len(r) for r in block.table)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for c in range(cols):
        table.columns[c].width = width // cols
    for r in range(rows):
        table.rows[r].height = height // rows
    base_style = block.text_style or {}
    for r, row in enumerate(block.table):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell_info = None
            if block.table_cells and r < len(block.table_cells):
                row_cells = block.table_cells[r]
                if c < len(row_cells):
                    cell_info = row_cells[c]
            text_value = value
            if cell_info and cell_info.text:
                text_value = cell_info.text
            cell.text = text_value
            tf = cell.text_frame
            tf.word_wrap = True
            paragraph = tf.paragraphs[0]
            style = dict(base_style)
            if cell_info:
                cell_style = cell_info.text_style or {}
                style.update({k: v for k, v in cell_style.items() if v is not None})
            apply_paragraph_alignment(paragraph, style.get("align"))
            font = paragraph.font
            size = style.get("font_size", DEFAULT_FONT_SIZE)
            font.size = Pt(px_to_pt(size))
            if style.get("bold") is not None:
                font.bold = bool(style["bold"])
            if style.get("italic") is not None:
                font.italic = bool(style["italic"])
            color = style.get("color")
            rgb = css_color_to_rgb_tuple(color)
            if rgb:
                font.color.rgb = RGBColor(*rgb)
            if cell_info:
                bg = cell_info.background_color
                if bg and not css_is_transparent(bg):
                    rgb_bg = css_color_to_rgb_tuple(bg)
                    if rgb_bg:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*rgb_bg)
                else:
                    cell.fill.background()
                if cell_info.border_color:
                    apply_cell_border(cell, cell_info.border_color, cell_info.border_width or 1.0)


def resolve_image_size(block: Block) -> Tuple[float, float]:
    width = block.layout.width or 640
    height = block.layout.height or 360
    if block.image_path and Image:
        try:
            with Image.open(block.image_path) as img:
                natural_w, natural_h = img.size
            if block.layout.width and block.layout.height:
                return block.layout.width, block.layout.height
            if block.layout.width and not block.layout.height:
                ratio = block.layout.width / natural_w
                return block.layout.width, natural_h * ratio
            if block.layout.height and not block.layout.width:
                ratio = block.layout.height / natural_h
                return natural_w * ratio, block.layout.height
            return natural_w, natural_h
        except Exception:
            return width, height
    return width, height


def add_image_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    width_px, height_px = resolve_image_size(block)
    left = position_to_emu(block.layout.left or DEFAULT_PADDING_X, "x", prs, slide_model)
    top = position_to_emu(block.layout.top or DEFAULT_PADDING_Y, "y", prs, slide_model)
    width = length_to_emu(width_px, "x", prs, slide_model)
    height = length_to_emu(height_px, "y", prs, slide_model)
    if block.image_path and block.image_path.exists():
        slide.shapes.add_picture(str(block.image_path), left, top, width=width, height=height)
    else:
        placeholder = slide.shapes.add_textbox(left, top, width, height)
        tf = placeholder.text_frame
        tf.text = block.image_alt or "[画像]"
        paragraph = tf.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(px_to_pt(24))
        placeholder.line.color.rgb = RGBColor(200, 200, 200)
        placeholder.line.width = Pt(px_to_pt(2))


def add_shape_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    left = position_to_emu(block.layout.left or DEFAULT_PADDING_X, "x", prs, slide_model)
    top = position_to_emu(block.layout.top or DEFAULT_PADDING_Y, "y", prs, slide_model)
    width = length_to_emu(block.layout.width or 240, "x", prs, slide_model)
    height = length_to_emu(block.layout.height or 120, "y", prs, slide_model)
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if block.shape_style.get("border_radius")
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    fill_color = css_color_to_rgb_tuple(block.shape_style.get("fill_color"))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*fill_color)
    else:
        shape.fill.background()
    border_color = css_color_to_rgb_tuple(block.shape_style.get("border_color"))
    if border_color:
        shape.line.color.rgb = RGBColor(*border_color)
        width_px = block.shape_style.get("border_width") or 2
        shape.line.width = Pt(px_to_pt(width_px))
    else:
        shape.line.fill.background()


def add_polyline_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    points = block.vector_data.get("points") or []
    stroke_color = block.vector_data.get("stroke")
    stroke_width = float(block.vector_data.get("stroke_width") or 2.0)
    fill_color = block.vector_data.get("fill")
    closed = bool(block.vector_data.get("closed"))
    if len(points) < 2:
        return
    if closed and fill_color and fill_color.lower() != "none":
        # approximate closed polygon via Freeform
        builder = slide.shapes.build_freeform(
            position_to_emu(points[0]["x"], "x", prs, slide_model),
            position_to_emu(points[0]["y"], "y", prs, slide_model),
        )
        segments = [
            (
                position_to_emu(pt["x"], "x", prs, slide_model),
                position_to_emu(pt["y"], "y", prs, slide_model),
            )
            for pt in points[1:]
        ]
        builder.add_line_segments(segments, close=True)
        shape = builder.convert_to_shape()
        rgb_fill = css_color_to_rgb_tuple(fill_color)
        if rgb_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb_fill)
        else:
            shape.fill.background()
        if stroke_color and stroke_color.lower() != "none":
            rgb = css_color_to_rgb_tuple(stroke_color)
            if rgb:
                shape.line.color.rgb = RGBColor(*rgb)
                shape.line.width = Pt(px_to_pt(stroke_width))
    else:
        for i in range(len(points) - 1):
            p1 = points[i]
            p2 = points[i + 1]
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                position_to_emu(p1["x"], "x", prs, slide_model),
                position_to_emu(p1["y"], "y", prs, slide_model),
                position_to_emu(p2["x"], "x", prs, slide_model),
                position_to_emu(p2["y"], "y", prs, slide_model),
            )
            if stroke_color and stroke_color.lower() != "none":
                rgb = css_color_to_rgb_tuple(stroke_color)
                if rgb:
                    connector.line.color.rgb = RGBColor(*rgb)
            connector.line.width = Pt(px_to_pt(stroke_width))


def add_circle_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    circle = block.vector_data or {}
    cx = float(circle.get("cx") or 0.0)
    cy = float(circle.get("cy") or 0.0)
    r = float(circle.get("r") or 0.0)
    if r <= 0:
        return
    left = position_to_emu(cx - r, "x", prs, slide_model)
    top = position_to_emu(cy - r, "y", prs, slide_model)
    size = length_to_emu(r * 2, "x", prs, slide_model)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, size, size)
    fill = circle.get("fill")
    stroke = circle.get("stroke")
    stroke_width = float(circle.get("strokeWidth") or 1.0)
    if fill and fill.lower() != "none":
        rgb = css_color_to_rgb_tuple(fill)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
    else:
        shape.fill.background()
    if stroke and stroke.lower() != "none":
        rgb = css_color_to_rgb_tuple(stroke)
        if rgb:
            shape.line.color.rgb = RGBColor(*rgb)
            shape.line.width = Pt(px_to_pt(stroke_width))


def add_ellipse_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    ellipse = block.vector_data or {}
    cx = float(ellipse.get("cx") or 0.0)
    cy = float(ellipse.get("cy") or 0.0)
    rx = float(ellipse.get("rx") or 0.0)
    ry = float(ellipse.get("ry") or 0.0)
    if rx <= 0 or ry <= 0:
        return
    left = position_to_emu(cx - rx, "x", prs, slide_model)
    top = position_to_emu(cy - ry, "y", prs, slide_model)
    width = length_to_emu(rx * 2, "x", prs, slide_model)
    height = length_to_emu(ry * 2, "y", prs, slide_model)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
    fill = ellipse.get("fill")
    stroke = ellipse.get("stroke")
    stroke_width = float(ellipse.get("strokeWidth") or 1.0)
    if fill and fill.lower() != "none":
        rgb = css_color_to_rgb_tuple(fill)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
    else:
        shape.fill.background()
    if stroke and stroke.lower() != "none":
        rgb = css_color_to_rgb_tuple(stroke)
        if rgb:
            shape.line.color.rgb = RGBColor(*rgb)
            shape.line.width = Pt(px_to_pt(stroke_width))


def add_conic_gradient_block(slide, block: Block, prs: Presentation, slide_model: SlideModel) -> None:
    gradient = block.vector_data or {}
    segments = gradient.get("segments") or []
    if not segments:
        return
    cx = gradient.get("cx")
    cy = gradient.get("cy")
    radius = gradient.get("radius")
    if not radius:
        width = block.layout.width or 240
        height = block.layout.height or width
        radius = min(width, height) / 2
    if cx is None or cy is None:
        left = block.layout.left or DEFAULT_PADDING_X
        top = block.layout.top or DEFAULT_PADDING_Y
        width = block.layout.width or radius * 2
        height = block.layout.height or radius * 2
        cx = left + width / 2
        cy = top + height / 2
    for seg in segments:
        color = seg.get("color")
        start_deg = float(seg.get("startDeg") or 0.0)
        end_deg = float(seg.get("endDeg") or 0.0)
        if end_deg <= start_deg or not color:
            continue
        span = end_deg - start_deg
        steps = max(3, int(span / 8))
        def css_deg_to_point(deg):
            rad = math.radians(90 - deg)
            return (
                cx + radius * math.cos(rad),
                cy + radius * math.sin(rad),
            )
        builder = slide.shapes.build_freeform(
            position_to_emu(cx, "x", prs, slide_model),
            position_to_emu(cy, "y", prs, slide_model),
        )
        start_point = css_deg_to_point(start_deg)
        segments = [
            (
                position_to_emu(start_point[0], "x", prs, slide_model),
                position_to_emu(start_point[1], "y", prs, slide_model),
            )
        ]
        for i in range(1, steps + 1):
            angle = start_deg + (span * i / steps)
            point = css_deg_to_point(angle)
            segments.append(
                (
                    position_to_emu(point[0], "x", prs, slide_model),
                    position_to_emu(point[1], "y", prs, slide_model),
                )
            )
        builder.add_line_segments(segments, close=True)
        shape = builder.convert_to_shape()
        rgb = css_color_to_rgb_tuple(color)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*rgb)
        else:
            shape.fill.background()
        shape.line.fill.background()


def slide_model_to_pptx(slides: List[SlideModel], output_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_model in slides:
        slide = prs.slides.add_slide(blank)
        if slide_model.background_color:
            rgb = css_color_to_rgb_tuple(slide_model.background_color)
            if rgb:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*rgb)
        ensure_slide_transform(slide_model, prs)
        ordered_blocks = sorted(slide_model.blocks, key=lambda b: (b.z_index, b.order))
        for block in ordered_blocks:
            try:
                if block.kind == "text":
                    add_text_block(slide, block, prs, slide_model)
                elif block.kind == "list":
                    add_list_block(slide, block, prs, slide_model)
                elif block.kind == "table":
                    add_table_block(slide, block, prs, slide_model)
                elif block.kind == "image":
                    add_image_block(slide, block, prs, slide_model)
                elif block.kind == "shape":
                    add_shape_block(slide, block, prs, slide_model)
                elif block.kind == "vector_polyline":
                    add_polyline_block(slide, block, prs, slide_model)
                elif block.kind == "vector_circle":
                    add_circle_block(slide, block, prs, slide_model)
                elif block.kind == "vector_ellipse":
                    add_ellipse_block(slide, block, prs, slide_model)
                elif block.kind == "conic-gradient":
                    add_conic_gradient_block(slide, block, prs, slide_model)
            except Exception as exc:
                print(f"[warn] スライド要素の描画に失敗しました: {exc}", file=sys.stderr)
    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="HTML を編集可能な PPTX に変換します。")
    parser.add_argument("input_html", help="入力 HTML ファイルパス")
    parser.add_argument("output_pptx", help="出力 PPTX パス")
    parser.add_argument(
        "--selector",
        default=".slide, .slide-container, [data-slide]",
        help="スライド要素として扱う CSS セレクタ（該当しない場合は <body> を1枚として扱う）",
    )
    args = parser.parse_args()

    input_path = Path(args.input_html)
    if not input_path.exists():
        print(f"入力ファイルが存在しません: {input_path}", file=sys.stderr)
        sys.exit(1)

    try:
        slides = parse_html_static(str(input_path), selector=args.selector)
        slide_model_to_pptx(slides, args.output_pptx)
    except Exception as exc:
        print(f"変換に失敗しました: {exc}", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
